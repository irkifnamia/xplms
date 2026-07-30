"""XPLMS — Experience Point Learning Management System."""

from __future__ import annotations

import base64
import html
import io
import hashlib
import json
import os
import re
from datetime import datetime
from pathlib import Path
from typing import Any
from zoneinfo import ZoneInfo

import bcrypt
import pandas as pd
import plotly.express as px
import streamlit as st
from PIL import Image
from pydantic import BaseModel, Field

try:
    from docx import Document
    from openai import OpenAI
    from pptx import Presentation
    from pypdf import PdfReader
except ImportError:  # Features remain unavailable until requirements are installed.
    Document = OpenAI = Presentation = PdfReader = None

try:
    from supabase import Client, create_client
except ImportError:  # pragma: no cover
    Client = Any
    create_client = None


class GeneratedQuestion(BaseModel):
    question: str
    options: list[str] = Field(min_length=4, max_length=4)
    correct_index: int = Field(ge=0, le=3)
    explanation: str


class GeneratedQuiz(BaseModel):
    title: str
    questions: list[GeneratedQuestion] = Field(min_length=1, max_length=20)


LOGO_PATH = Path(__file__).parent / "assets" / "xplms-logo.png"
AFJ_LOGO_PATH = Path(__file__).parent / "assets" / "afj.jpeg"
LOGO_IMAGE = Image.open(LOGO_PATH)
LOGO_DATA = base64.b64encode(LOGO_PATH.read_bytes()).decode("ascii")

st.set_page_config(
    page_title="XPLMS",
    page_icon=LOGO_IMAGE,
    layout="wide",
    initial_sidebar_state="expanded",
)

st.markdown(
    """
    <style>
    [data-testid="stToolbar"],
    [data-testid="stToolbarActions"],
    [data-testid="stStatusWidget"],
    [data-testid="stDecoration"],
    [data-testid="stDeployButton"],
    .stDeployButton,
    #MainMenu {
        display: none !important;
        visibility: hidden !important;
        height: 0 !important;
    }

    [data-testid="stHeader"] {
        display: none !important;
        visibility: hidden !important;
        height: 0 !important;
        min-height: 0 !important;
        background: transparent !important;
        pointer-events: none !important;
    }

    [data-testid="stExpandSidebarButton"],
    [data-testid="collapsedControl"] {
        display: flex !important;
        visibility: visible !important;
        position: fixed !important;
        top: 0.75rem !important;
        left: 0.75rem !important;
        z-index: 1000000 !important;
        pointer-events: auto !important;
    }

    [data-testid="stExpandSidebarButton"],
    [data-testid="collapsedControl"] button {
        width: 2.75rem !important;
        height: 2.75rem !important;
        min-height: 2.75rem !important;
        padding: 0.45rem !important;
        border: 1px solid #c4d1df !important;
        border-radius: 12px !important;
        background: #ffffff !important;
        box-shadow: 0 4px 14px rgba(23,35,58,.14) !important;
    }

    [data-testid="stAppViewContainer"] .block-container {
        padding-top: 0.5rem !important;
    }
    </style>
    """,
    unsafe_allow_html=True,
)

BRAND = "#1F5DAA"
CSS = """
<style>
@import url('https://fonts.googleapis.com/css2?family=DM+Sans:wght@400;500;600;700&family=Manrope:wght@600;700;800&display=swap');
:root { --ink:#17233a; --muted:#5f6b7a; --brand:#1f5daa; --brand-dark:#154b8c; --cyan:#12a9c4; --gold:#f2a900; --coral:#ef6a67; --soft:#f3f7fc; --line:#d3dce6; }
html, body, [class*="css"] { font-family:'DM Sans',sans-serif; color:var(--ink); }
h1,h2,h3 { font-family:'Manrope',sans-serif !important; letter-spacing:-.035em; }
[data-testid="stAppViewContainer"] {
  background:
    radial-gradient(circle at 92% 8%,rgba(18,169,196,.10),transparent 24rem),
    radial-gradient(circle at 8% 92%,rgba(242,169,0,.09),transparent 22rem),
    #f3f7fc;
}
[data-testid="stMainBlockContainer"] p,
[data-testid="stMainBlockContainer"] label,
[data-testid="stMainBlockContainer"] [data-testid="stMarkdownContainer"] { color:var(--ink); }
[data-testid="stMainBlockContainer"] [data-testid="stCaptionContainer"] { color:var(--muted); }
[data-testid="stSidebar"] {
  background:linear-gradient(180deg,#ffffff 0%,#f1f7ff 68%,#eefaf9 100%);
  border-right:1px solid #cbd9e8;
}
[data-testid="stSidebar"] [data-testid="stMarkdownContainer"],
[data-testid="stSidebar"] [data-testid="stMarkdownContainer"] p,
[data-testid="stSidebar"] label,
[data-testid="stSidebar"] .stCaption { color:#252525 !important; }
[data-testid="stSidebar"] .stRadio label {
  padding:.48rem .55rem; font-weight:600; border-radius:9px; color:#252525 !important;
}
[data-testid="stSidebar"] .stRadio label:hover { background:#dcecff; }
[data-testid="stSidebar"] .stRadio label:has(input:checked) {
  background:linear-gradient(90deg,#dcecff,#e5f8f7);
  box-shadow:inset 4px 0 0 var(--brand);
}
[data-testid="stSidebar"] .stRadio label p,
[data-testid="stSidebar"] .stRadio label span { color:#252525 !important; }
[data-testid="stSidebar"] hr { border-color:#d3dce6; }
[data-testid="stSidebar"] [data-baseweb="select"] > div {
  background:#fff !important; border-color:#aaa !important; color:#171717 !important;
}
[data-testid="stSidebar"] [data-baseweb="select"] span { color:#171717 !important; }
[data-testid="stSidebar"] [data-baseweb="select"] svg { fill:#171717; }
.block-container { max-width:1440px; padding:1.8rem 2.2rem 4rem; }
.topline { color:var(--brand); font-weight:700; text-transform:uppercase; letter-spacing:.12em; font-size:.72rem; }
.hero { background:linear-gradient(120deg,#fff 0%,#f4f9ff 58%,#e9f9f7 100%); color:#17233a;
  border:1px solid #dbe3ec; border-left:6px solid var(--brand); border-radius:18px;
  padding:26px 30px; margin-bottom:20px; box-shadow:0 10px 28px rgba(31,93,170,.09); position:relative; overflow:hidden; }
.hero:after { content:""; position:absolute; width:120px; height:120px; border-radius:50%; right:-35px; top:-45px; background:rgba(242,169,0,.16); }
.hero-inner { display:flex; align-items:center; justify-content:space-between; gap:20px; position:relative; z-index:1; }
.hero-brand { flex:0 0 auto; text-align:center; min-width:105px; }
.hero-brand img { display:block; width:72px; height:72px; object-fit:contain; margin:0 auto 3px; }
.hero-identity { color:#40516a; font-size:.68rem; font-weight:800; letter-spacing:.055em; white-space:nowrap; }
.hero h1 { color:#17233a; margin:.2rem 0 .35rem; font-size:2.15rem; }
.hero p { color:#536174 !important; margin:0; }
.metric-card,.panel { background:linear-gradient(145deg,#fff,#f8fbff); border:1px solid var(--line); border-top:4px solid var(--brand); border-radius:18px; padding:18px 20px;
  box-shadow:0 4px 14px rgba(0,0,0,.045); }
.metric-card:nth-child(3n+2) { border-top-color:var(--cyan); }
.metric-card:nth-child(3n) { border-top-color:var(--gold); }
.metric-label { color:var(--muted); font-size:.79rem; font-weight:700; text-transform:uppercase; letter-spacing:.06em; }
.metric-value { color:#17233a; font-family:'Manrope'; font-size:1.7rem; font-weight:800; margin:.15rem 0; }
.metric-note { color:var(--muted); font-size:.82rem; }
.result-grid { display:grid; grid-template-columns:repeat(3,minmax(0,1fr)); gap:14px; }
.result-card { background:#fff; border:1px solid var(--line); border-top:4px solid var(--brand);
  border-radius:16px; padding:17px 18px; box-shadow:0 4px 14px rgba(31,93,170,.07); }
.result-card:nth-child(3n+2) { border-top-color:var(--cyan); }
.result-card:nth-child(3n) { border-top-color:var(--gold); }
.result-title { color:var(--ink); font-weight:800; font-size:.86rem; letter-spacing:.035em; }
.result-mark { color:var(--brand-dark); font-family:'Manrope'; font-size:1.8rem; font-weight:800; margin:.3rem 0; }
.result-meta { color:var(--muted); font-size:.78rem; text-transform:uppercase; letter-spacing:.04em; }
.pill { display:inline-block; border-radius:999px; padding:4px 9px; font-size:.73rem; font-weight:700;
  color:#174c87; background:#e8f1fc; border:1px solid #b8d2ef; }
.brand { font-family:'Manrope'; font-size:1.35rem; font-weight:800; letter-spacing:-.04em; color:#17233a; }
.brandmark { display:inline-grid; place-items:center; width:29px; height:29px; margin-right:8px; border-radius:9px;
  background:var(--brand); color:white; }
.userbox { background:linear-gradient(120deg,#e7f1ff,#e6f8f6); border:1px solid #bdd6ed; color:#17233a; padding:12px; border-radius:14px; margin:14px 0 8px; }
.userbox b { color:#17233a !important; }
.userbox small { color:#5f6b7a !important; }
.empty { text-align:center; padding:44px 20px; color:var(--muted); background:white; border:1px dashed #ccd0dc; border-radius:18px; }
.stButton > button, .stDownloadButton > button { border-radius:10px; font-weight:700; min-height:2.65rem; color:#171717; border-color:#bdbdb9; background:#fff; }
.stButton > button p, .stDownloadButton > button p { color:#171717 !important; }
.stButton > button:hover, .stDownloadButton > button:hover { color:var(--brand); border-color:var(--brand); }
.stButton > button[kind="primary"] { background:var(--brand) !important; border-color:var(--brand) !important; color:white !important; }
.stButton > button[kind="primary"] p { color:white !important; }
.stButton > button[kind="primary"]:hover { background:var(--brand-dark) !important; border-color:var(--brand-dark) !important; color:white !important; }

/* Segmented filters: white unselected tabs, blue selected tab, readable labels. */
[data-testid="stSegmentedControl"] > div { background:white !important; border:1px solid #b8b8b4 !important; }
[data-testid="stSegmentedControl"] button,
[data-testid="stBaseButton-segmented_control"] {
  background:white !important; border-color:#d4d4d0 !important; color:#171717 !important;
}
[data-testid="stSegmentedControl"] button:hover,
[data-testid="stBaseButton-segmented_control"]:hover {
  background:#edf4fc !important; color:#154b8c !important;
}
[data-testid="stSegmentedControl"] button p,
[data-testid="stSegmentedControl"] button span,
[data-testid="stBaseButton-segmented_control"] p,
[data-testid="stBaseButton-segmented_control"] span { color:#171717 !important; }
[data-testid="stSegmentedControl"] button[aria-pressed="true"],
[data-testid="stBaseButton-segmented_control"][aria-pressed="true"] {
  background:var(--brand) !important; border-color:var(--brand) !important;
}
[data-testid="stSegmentedControl"] button[aria-pressed="true"] p,
[data-testid="stSegmentedControl"] button[aria-pressed="true"] span,
[data-testid="stBaseButton-segmented_control"][aria-pressed="true"] p,
[data-testid="stBaseButton-segmented_control"][aria-pressed="true"] span { color:white !important; }

/* Form controls must stay light even when the OS or Streamlit defaults are dark. */
[data-testid="stFileUploader"] { border-radius:14px; background:white !important; color:#171717 !important; }
[data-testid="stFileUploaderDropzone"],
[data-testid="stFileUploaderDropzone"] > div {
  background:#fafafa !important; border-color:#b8b8b4 !important; color:#171717 !important;
}
[data-testid="stFileUploaderDropzone"] p,
[data-testid="stFileUploaderDropzone"] span,
[data-testid="stFileUploaderDropzone"] small { color:#363636 !important; }
[data-testid="stFileUploaderDropzone"] button {
  background:white !important; border:1px solid #9b9b97 !important; color:#171717 !important;
}
[data-testid="stFileUploaderDropzone"] button p { color:#171717 !important; }
[data-baseweb="input"] > div, [data-baseweb="textarea"] > div,
[data-testid="stTextInput"] > div > div, [data-testid="stTextArea"] > div > div {
  background:white !important; color:#171717 !important;
}
[data-baseweb="input"] input, [data-baseweb="textarea"] textarea,
[data-testid="stTextInput"] input, [data-testid="stTextArea"] textarea {
  background:white !important; color:#171717 !important; -webkit-text-fill-color:#171717 !important;
  border-color:#aaa !important; caret-color:var(--brand);
}
[data-baseweb="input"] input::placeholder, [data-baseweb="textarea"] textarea::placeholder { color:#737373 !important; opacity:1; }
[data-baseweb="tab-list"] { background:#e8f0f9; border-radius:12px; padding:4px; gap:3px; }
[data-testid="stAlert"] { color:#171717; }
[data-testid="stAlert"] p, [data-testid="stAlert"] span { color:#171717 !important; }
[aria-disabled="true"], button:disabled { opacity:1 !important; }
button:disabled, button:disabled p, button:disabled span {
  background:#efefed !important; color:#686868 !important; border-color:#d3d3cf !important;
}
[data-testid="stDataFrame"] { border:1px solid var(--line); border-radius:14px; overflow:hidden; }

/* Final accessibility layer: never rely on hover to reveal wording. */
[data-testid="stHeader"],
[data-testid="stToolbar"],
[data-testid="stDecoration"] { background:#fff !important; color:#171717 !important; }
[data-testid="stHeader"] * { color:#171717 !important; }
[data-testid="stAppViewContainer"],
[data-testid="stMain"],
[data-testid="stMainBlockContainer"] { color:#171717 !important; }
[data-testid="stMainBlockContainer"] label,
[data-testid="stMainBlockContainer"] label *,
[data-testid="stMainBlockContainer"] p,
[data-testid="stMainBlockContainer"] span { text-shadow:none !important; }
[data-testid="stSidebar"] [data-testid="stCaptionContainer"],
[data-testid="stSidebar"] [data-testid="stCaptionContainer"] *,
[data-testid="stSidebar"] small { color:#5b5b57 !important; opacity:1 !important; }
[data-testid="stImage"] {
  width:100% !important;
  display:flex !important;
  justify-content:center !important;
  margin-left:auto !important;
  margin-right:auto !important;
  text-align:center !important;
}
[data-testid="stImage"] img {
  display:block !important;
  margin-left:auto !important;
  margin-right:auto !important;
}

/* A collapsed preference must never make desktop navigation disappear. */
@media(min-width:701px){
  [class*="st-key-mobile_navigation_"] {
    display:none !important;
  }
  [data-testid="stVerticalBlockBorderWrapper"]:has(.mobile-nav-marker) {
    display:none !important;
  }
  section[data-testid="stSidebar"] {
    display:block !important;
    visibility:visible !important;
    min-width:280px !important;
    width:280px !important;
    transform:translateX(0) !important;
  }
}

/* Native buttons, form buttons and segmented button groups. */
button[kind="secondary"],
button[kind="tertiary"],
button[kind="segmented_control"],
[data-testid*="BaseButton-secondary"],
[data-testid*="BaseButton-tertiary"],
[data-testid*="BaseButton-segmented"] {
  background:#fff !important;
  color:#171717 !important;
  border-color:#aaa9a4 !important;
}
button[kind="secondary"] *,
button[kind="tertiary"] *,
button[kind="segmented_control"] *,
[data-testid*="BaseButton-secondary"] *,
[data-testid*="BaseButton-tertiary"] *,
[data-testid*="BaseButton-segmented"] * { color:#171717 !important; opacity:1 !important; }
button[kind="primary"],
[data-testid*="BaseButton-primary"] {
  background:#1f5daa !important; color:#fff !important; border-color:#1f5daa !important;
}
button[kind="primary"] *,
[data-testid*="BaseButton-primary"] * { color:#fff !important; opacity:1 !important; }
button[kind="segmented_control"][aria-pressed="true"],
button[kind="segmented_control"][aria-checked="true"],
[data-testid*="BaseButton-segmented"][aria-pressed="true"],
[data-testid*="BaseButton-segmented"][aria-checked="true"] {
  background:#1f5daa !important; border-color:#1f5daa !important;
}
button[kind="segmented_control"][aria-pressed="true"] *,
button[kind="segmented_control"][aria-checked="true"] *,
[data-testid*="BaseButton-segmented"][aria-pressed="true"] *,
[data-testid*="BaseButton-segmented"][aria-checked="true"] * { color:#fff !important; }

/* Tabs and menu wording. */
[data-baseweb="tab-list"] { background:#e8f0f9 !important; border:1px solid #cfdeed; }
[data-baseweb="tab"] { color:#2b2b29 !important; opacity:1 !important; }
[data-baseweb="tab"] * { color:#2b2b29 !important; opacity:1 !important; }
[data-baseweb="tab"][aria-selected="true"],
[data-baseweb="tab"][aria-selected="true"] {
  background:#fff !important; border-radius:9px; box-shadow:0 2px 8px rgba(31,93,170,.12);
}
[data-baseweb="tab"][aria-selected="true"],
[data-baseweb="tab"][aria-selected="true"] * { color:#154b8c !important; font-weight:700; }

/* Every field has a light surface and a dark value/placeholder. */
input, textarea, [contenteditable="true"],
[data-baseweb="select"] > div {
  background:#fff !important; color:#171717 !important; -webkit-text-fill-color:#171717 !important;
}
input::placeholder, textarea::placeholder { color:#696965 !important; opacity:1 !important; }
[data-baseweb="select"] *, [role="option"], [role="option"] * {
  color:#171717 !important; opacity:1 !important;
}
[role="listbox"], [role="option"] { background:#fff !important; }

@media(max-width:700px){
  /* Keep Streamlit chrome hidden, but restore the mobile sidebar opener for Admin. */
  [data-testid="stHeader"] {
    display:none !important;
    visibility:hidden !important;
    height:0 !important;
    min-height:0 !important;
    background:transparent !important;
  }
  [data-testid="stHeader"] [data-testid="stToolbar"],
  [data-testid="stHeader"] [data-testid="stToolbarActions"],
  [data-testid="stHeader"] [data-testid="stStatusWidget"],
  [data-testid="stHeader"] [data-testid="stDeployButton"] {
    display:none !important;
  }
  [data-testid="collapsedControl"],
  [data-testid="stExpandSidebarButton"] {
    display:flex !important;
    visibility:visible !important;
    position:fixed !important;
    top:.7rem !important;
    left:.7rem !important;
    z-index:1000000 !important;
  }
  [data-testid="collapsedControl"] button,
  [data-testid="stExpandSidebarButton"] {
    width:2.75rem !important;
    height:2.75rem !important;
    min-height:2.75rem !important;
    padding:.45rem !important;
    border:1px solid #c4d1df !important;
    border-radius:12px !important;
    background:#fff !important;
    box-shadow:0 4px 14px rgba(23,35,58,.14) !important;
  }
  [data-testid="collapsedControl"] button *,
  [data-testid="collapsedControl"] svg,
  [data-testid="stExpandSidebarButton"] *,
  [data-testid="stExpandSidebarButton"] svg {
    color:#17233a !important;
    fill:#17233a !important;
  }
  body:has(.student-sidebar-marker) .block-container {
    padding:.35rem 0.72rem 4.5rem !important;
  }
  body:has(.student-sidebar-marker) .hero {
    padding:15px 16px;
    margin-bottom:12px;
    border-left-width:4px;
    border-radius:13px;
  }
  body:has(.student-sidebar-marker) .hero h1 {
    font-size:1.42rem;
    line-height:1.18;
  }
  .hero-inner { gap:10px; }
  .hero-brand { min-width:82px; }
  .hero-brand img { width:54px; height:54px; }
  .hero-identity { font-size:.58rem; max-width:110px; overflow:hidden; text-overflow:ellipsis; }
  body:has(.student-sidebar-marker) .hero p {
    font-size:.88rem;
    line-height:1.4;
  }
  body:has(.student-sidebar-marker) .topline { font-size:.64rem; }
  body:has(.student-sidebar-marker) .metric-card,
  body:has(.student-sidebar-marker) .panel {
    padding:13px 14px;
    border-radius:13px;
  }
  body:has(.student-sidebar-marker) .metric-value { font-size:1.32rem; }
  body:has(.student-sidebar-marker) [data-testid="stVerticalBlockBorderWrapper"] {
    border-radius:13px !important;
  }
  body:has(.student-sidebar-marker) [data-testid="stPlotlyChart"] {
    background:#fff;
    border:1px solid var(--line);
    border-radius:13px;
    overflow:hidden;
  }
  body:has(.student-sidebar-marker) [data-testid="stDataFrame"] {
    font-size:.78rem;
  }
  .result-grid { grid-template-columns:1fr; gap:10px; }
  body:has(.student-sidebar-marker) button {
    min-height:2.75rem;
    font-size:.88rem;
  }
  body:has(.student-sidebar-marker) [data-testid="stSegmentedControl"] {
    overflow-x:auto;
    scrollbar-width:none;
  }
  body:has(.student-sidebar-marker) [data-testid="stSegmentedControl"] > div {
    min-width:max-content;
  }

  .block-container { padding:.35rem .9rem 4.5rem !important; }
  [data-testid="stVerticalBlockBorderWrapper"]:has(.mobile-nav-marker) {
    display:flex !important;
    margin-bottom:.25rem;
  }
  [class*="st-key-mobile_navigation_"] {
    display:block !important;
    margin-top:0 !important;
    padding-top:0 !important;
  }
  [class*="st-key-mobile_navigation_"] [data-baseweb="select"],
  [class*="st-key-mobile_navigation_"] [data-baseweb="select"] > div,
  [class*="st-key-mobile_navigation_"] [data-baseweb="select"] > div > div {
    background-color:#17834d !important;
    border-color:#17834d !important;
    color:#fff !important;
  }
  [class*="st-key-mobile_navigation_"] [data-baseweb="select"] span,
  [class*="st-key-mobile_navigation_"] [data-baseweb="select"] input {
    color:#fff !important;
    -webkit-text-fill-color:#fff !important;
  }
  [class*="st-key-mobile_navigation_"] [data-baseweb="select"] svg {
    fill:#fff !important;
  }
  [data-testid="stVerticalBlockBorderWrapper"]:has(.mobile-nav-marker)
  [data-testid="stHorizontalBlock"] {
    display:flex !important;
    flex-direction:row !important;
    flex-wrap:nowrap !important;
    gap:.55rem !important;
  }
  [data-testid="stVerticalBlockBorderWrapper"]:has(.mobile-nav-marker)
  [data-testid="stColumn"] {
    min-width:0 !important;
  }
  [class*="st-key-mobile_sign_out_"] button {
    background:#c62828 !important;
    border-color:#c62828 !important;
    color:#fff !important;
  }
  [class*="st-key-mobile_sign_out_"] button *,
  [class*="st-key-mobile_sign_out_"] button p {
    color:#fff !important;
  }
  [class*="st-key-mobile_sign_out_"] button:hover {
    background:#a91f1f !important;
    border-color:#a91f1f !important;
  }
  [data-testid="stVerticalBlockBorderWrapper"]:has(.mobile-nav-marker)
  [data-baseweb="select"],
  [data-testid="stVerticalBlockBorderWrapper"]:has(.mobile-nav-marker)
  [data-baseweb="select"] > div,
  [data-testid="stVerticalBlockBorderWrapper"]:has(.mobile-nav-marker)
  [data-baseweb="select"] > div > div {
    background-color:#17834d !important;
    border-color:#17834d !important;
    color:#fff !important;
    box-shadow:none !important;
  }
  [data-testid="stVerticalBlockBorderWrapper"]:has(.mobile-nav-marker)
  [data-baseweb="select"] span,
  [data-testid="stVerticalBlockBorderWrapper"]:has(.mobile-nav-marker)
  [data-baseweb="select"] input {
    color:#fff !important;
    -webkit-text-fill-color:#fff !important;
    background:transparent !important;
  }
  [data-testid="stVerticalBlockBorderWrapper"]:has(.mobile-nav-marker)
  [data-baseweb="select"] svg {
    fill:#fff !important;
  }
  .hero { padding:20px; border-radius:17px; }
  .hero h1 { font-size:1.55rem; }
  .metric-value { font-size:1.4rem; }
}
</style>
"""
st.markdown(CSS, unsafe_allow_html=True)


DEMO_STUDENTS = pd.DataFrame(
    [
        {"student_id": "S24001", "name": "Alya Sofea", "programme": "Digital Technology", "cohort": "2024", "email": "alya@example.edu", "XPTEAM": "J1T2(A)"},
        {"student_id": "S24002", "name": "Haziq Amir", "programme": "Digital Technology", "cohort": "2024", "email": "haziq@example.edu", "XPTEAM": "J1T2(A)"},
        {"student_id": "S24003", "name": "Mei Lin", "programme": "Information Systems", "cohort": "2024", "email": "meilin@example.edu"},
        {"student_id": "S24004", "name": "Arjun Kumar", "programme": "Information Systems", "cohort": "2024", "email": "arjun@example.edu"},
        {"student_id": "S24005", "name": "Nur Iman", "programme": "Digital Technology", "cohort": "2024", "email": "iman@example.edu"},
    ]
)
DEMO_PROGRESS = pd.DataFrame(
    [
        {"student_id": "S24001", "assignment": 87, "pb": 82, "task": 94, "xp": 2460, "badge": 8},
        {"student_id": "S24002", "assignment": 81, "pb": 78, "task": 89, "xp": 2180, "badge": 7},
        {"student_id": "S24003", "assignment": 92, "pb": 91, "task": 86, "xp": 2720, "badge": 10},
        {"student_id": "S24004", "assignment": 74, "pb": 80, "task": 77, "xp": 1840, "badge": 5},
        {"student_id": "S24005", "assignment": 88, "pb": 85, "task": 91, "xp": 2530, "badge": 9},
    ]
)
EMPTY_MATERIALS = pd.DataFrame(
    columns=[
        "id", "title", "course", "chapter", "material_type", "type",
        "file_path", "uploaded_at",
    ]
)


def db() -> Client | None:
    if create_client is None:
        return None
    if "_supabase_client" in st.session_state:
        return st.session_state["_supabase_client"]
    try:
        service_key = st.secrets.get("SUPABASE_SECRET_KEY") or st.secrets.get(
            "SUPABASE_SERVICE_ROLE_KEY"
        )
        if not service_key:
            return None
        client = create_client(st.secrets["SUPABASE_URL"], service_key)
        st.session_state["_supabase_client"] = client
        return client
    except Exception:
        return None


def hash_password(password: str) -> str:
    return bcrypt.hashpw(password.encode("utf-8"), bcrypt.gensalt(rounds=12)).decode("utf-8")


def valid_password(password: str) -> bool:
    return (
        len(password) >= 10
        and any(c.isupper() for c in password)
        and any(c.islower() for c in password)
        and any(c.isdigit() for c in password)
    )


def is_demo() -> bool:
    """Keep preview sessions completely isolated from Supabase."""
    user = st.session_state.get("user", {})
    return user.get("id") == "demo"


def fetch_table(name: str, fallback: pd.DataFrame) -> tuple[pd.DataFrame, bool]:
    if is_demo():
        return fallback.copy(), False
    client = db()
    if client:
        try:
            rows = client.table(name).select("*").execute().data
            return (pd.DataFrame(rows), True) if rows else (fallback.copy(), True)
        except Exception:
            pass
    return fallback.copy(), False


def fetch_student_row(
    name: str, matric: Any, fallback: pd.DataFrame
) -> tuple[pd.DataFrame, bool]:
    """Fetch one student's row instead of transferring the whole table."""
    if is_demo():
        return fallback.head(1).copy(), False
    client = db()
    if client and matric:
        try:
            rows = (
                client.table(name).select("*")
                .eq("NO MATRIK", str(matric)).limit(1).execute().data
            )
            return pd.DataFrame(rows), True
        except Exception:
            pass
    return fallback.head(0).copy(), False


def upsert_rows(name: str, frame: pd.DataFrame) -> tuple[bool, str]:
    if is_demo():
        return False, "Demo mode is read-only. No Supabase data was changed."
    client = db()
    if not client:
        return False, "Supabase is unavailable. No data was changed."
    try:
        clean = frame.where(pd.notna(frame), None).to_dict("records")
        client.table(name).upsert(clean).execute()
        return True, f"{len(clean)} records saved."
    except Exception as exc:
        return False, f"Could not save records: {exc}"


def upload_file(bucket: str, path: str, data: bytes, content_type: str) -> tuple[bool, str]:
    if is_demo():
        return False, "Demo mode is read-only. No file was uploaded."
    client = db()
    if not client:
        return False, "Supabase storage is unavailable."
    try:
        client.storage.from_(bucket).upload(path, data, {"content-type": content_type, "upsert": "true"})
        return True, path
    except Exception as exc:
        return False, str(exc)


def extract_material_text(filename: str, data: bytes) -> str:
    extension = filename.lower().rsplit(".", 1)[-1]
    if extension == "pdf" and PdfReader:
        return "\n".join(page.extract_text() or "" for page in PdfReader(io.BytesIO(data)).pages)
    if extension == "docx" and Document:
        document = Document(io.BytesIO(data))
        return "\n".join(paragraph.text for paragraph in document.paragraphs)
    if extension == "pptx" and Presentation:
        presentation = Presentation(io.BytesIO(data))
        return "\n".join(
            shape.text
            for slide in presentation.slides
            for shape in slide.shapes
            if hasattr(shape, "text")
        )
    if extension in {"txt", "md", "csv"}:
        return data.decode("utf-8", errors="ignore")
    raise ValueError("Quiz generation supports PDF, DOCX, PPTX, TXT, MD and CSV materials.")


def generate_material_quiz(material: pd.Series, question_count: int) -> tuple[bool, str]:
    api_key = st.secrets.get("OPENAI_API_KEY")
    if not api_key or OpenAI is None:
        return False, "Add OPENAI_API_KEY to Streamlit secrets and install requirements to enable AI quiz generation."
    client = db()
    try:
        file_path = str(material["file_path"])
        data = client.storage.from_("materials").download(file_path)
        source_text = extract_material_text(file_path, data).strip()
        if len(source_text) < 200:
            return False, "The selected material does not contain enough readable text."
        source_text = source_text[:60_000]
        generated_questions: list[GeneratedQuestion] = []
        generated_title = f"{material['title']} quiz"
        openai_client = OpenAI(api_key=api_key)
        remaining = question_count
        batch_number = 1
        while remaining:
            batch_size = min(20, remaining)
            response = openai_client.responses.parse(
                model="gpt-5.6-luna",
                input=[
                    {
                        "role": "system",
                        "content": (
                            "Create a fair educational multiple-choice quiz using only the supplied "
                            "study material. Use four plausible options per question, one correct answer, "
                            "clear explanations, and no trick questions. Avoid repeating questions."
                        ),
                    },
                    {
                        "role": "user",
                        "content": (
                            f"Create exactly {batch_size} questions for batch {batch_number}. "
                            f"Do not repeat these existing questions: "
                            f"{[item.question for item in generated_questions]}\n\n"
                            f"Study material:\n{source_text}"
                        ),
                    },
                ],
                text_format=GeneratedQuiz,
            )
            generated = response.output_parsed
            if batch_number == 1 and generated.title:
                generated_title = generated.title
            generated_questions.extend(generated.questions)
            remaining -= batch_size
            batch_number += 1
        quiz_row = client.table("quizzes").insert({
            "material_id": int(material["id"]),
            "title": generated_title,
            "instructions": "Choose the best answer for each question.",
            "xp_reward": 25,
            "passing_score": 60,
            "status": "draft",
            "generated_by_ai": True,
            "created_by": st.session_state.user["id"],
        }).execute().data[0]
        questions = [
            {
                "quiz_id": quiz_row["id"],
                "position": index,
                "question": item.question,
                "options": item.options,
                "correct_index": item.correct_index,
                "explanation": item.explanation,
            }
            for index, item in enumerate(generated_questions, start=1)
        ]
        client.table("quiz_questions").insert(questions).execute()
        return True, f"Draft quiz created with {len(questions)} questions."
    except Exception as exc:
        return False, f"Quiz generation failed: {exc}"


def metric(label: str, value: str, note: str = "") -> None:
    note_html = f'<div class="metric-note">{note}</div>' if note else ""
    st.markdown(
        f'<div class="metric-card"><div class="metric-label">{label}</div>'
        f'<div class="metric-value">{value}</div>{note_html}</div>',
        unsafe_allow_html=True,
    )


def polish_chart(fig: Any) -> Any:
    """Keep Plotly typography readable regardless of the active Streamlit theme."""
    fig.update_layout(
        font=dict(family="DM Sans, sans-serif", color="#242424", size=13),
        title_font=dict(family="Manrope, sans-serif", color="#171717"),
        legend=dict(font=dict(color="#242424")),
        hoverlabel=dict(
            bgcolor="#ffffff",
            bordercolor="#1f5daa",
            font=dict(family="DM Sans, sans-serif", color="#171717", size=13),
        ),
    )
    fig.update_xaxes(
        tickfont=dict(color="#3f3f3f"),
        title_font=dict(color="#3f3f3f"),
        linecolor="#cfcfcb",
        gridcolor="#e9e9e6",
        zerolinecolor="#cfcfcb",
    )
    fig.update_yaxes(
        tickfont=dict(color="#3f3f3f"),
        title_font=dict(color="#3f3f3f"),
        linecolor="#cfcfcb",
        gridcolor="#e9e9e6",
        zerolinecolor="#cfcfcb",
    )
    return fig


def heading(eyebrow: str, title: str, copy: str = "") -> None:
    identity = html.escape(str(st.session_state.get("header_identity", "")))
    st.markdown(
        f'<div class="hero"><div class="hero-inner">'
        f'<h1>{html.escape(title.upper())}</h1>'
        f'<div class="hero-brand"><img src="data:image/png;base64,{LOGO_DATA}" '
        f'alt="XPLMS"><div class="hero-identity">{identity}</div></div>'
        f'</div></div>',
        unsafe_allow_html=True,
    )


def centered_image(path: Path, width: int, caption: str = "") -> None:
    """Render an image in a stable centred wrapper on desktop and mobile."""
    mime = "image/png" if path.suffix.lower() == ".png" else "image/jpeg"
    encoded = base64.b64encode(path.read_bytes()).decode("ascii")
    caption_html = (
        f'<div style="margin-bottom:7px;color:#5f6b7a;font-size:.68rem;'
        f'font-weight:800;letter-spacing:.14em">{caption}</div>'
        if caption
        else ""
    )
    st.markdown(
        f'<div style="width:100%;text-align:center">{caption_html}'
        f'<img src="data:{mime};base64,{encoded}" alt="" '
        f'style="display:block;width:{width}px;max-width:100%;height:auto;'
        f'margin:0 auto"></div>',
        unsafe_allow_html=True,
    )


def sync_navigation(source_key: str, active_key: str) -> None:
    st.session_state[active_key] = st.session_state[source_key]


def sidebar(role: str, name: str) -> str:
    with st.sidebar:
        if role == "Student":
            st.markdown('<span class="student-sidebar-marker"></span>', unsafe_allow_html=True)
        centered_image(LOGO_PATH, 112)
        st.markdown(
            f'<div class="userbox"><b>{name}</b><br><small>{role} workspace</small></div>',
            unsafe_allow_html=True,
        )
        if is_demo():
            roles = ["Student", "Admin"]
            selected_role = st.selectbox(
                "Demo workspace",
                roles,
                index=roles.index(role),
                help="Preview XPLMS from a different role.",
            )
            if selected_role != role:
                st.session_state.user.update(
                    {"name": f"Demo {selected_role}", "role": selected_role}
                )
                st.rerun()
        elif st.session_state.user.get("role") == "Admin":
            preview_roles = ["Admin", "Student"]
            selected_preview = st.selectbox(
                "Preview workspace",
                preview_roles,
                index=preview_roles.index(role),
                help="Preview another workspace without changing your administrator account.",
            )
            if selected_preview != role:
                st.session_state.admin_preview_role = selected_preview
                st.rerun()
        menus = {
            "Student": [
                "Profile", "Materials", "XP Journey", "SOP", "Leaderboard",
                "Results", "Request XP", "Quiz"
            ],
            "Admin": [
                "User access",
                "CRUD",
                "Material",
                "Quiz",
                "Award XP",
                "SOP",
                "Student record",
                "Analysis background",
                "Analysis results",
                "Analysis XP",
            ],
        }
        labels = {item: item.upper() for items in menus.values() for item in items}
        active_key = f"active_page_{role.lower()}"
        if st.session_state.get(active_key) not in menus[role]:
            st.session_state[active_key] = menus[role][0]
        sidebar_key = f"sidebar_page_{role.lower()}"
        if st.session_state.get(sidebar_key) not in menus[role]:
            st.session_state[sidebar_key] = st.session_state[active_key]
        st.radio(
            "Navigation",
            menus[role],
            format_func=lambda x: labels[x],
            label_visibility="collapsed",
            key=sidebar_key,
            on_change=sync_navigation,
            args=(sidebar_key, active_key),
        )
        st.divider()
        centered_image(AFJ_LOGO_PATH, 62, "EMPOWERED BY")
        if st.button("Sign out", use_container_width=True):
            st.session_state.clear()
            st.rerun()
    return st.session_state[active_key]


def mobile_navigation(role: str, current_page: str) -> str:
    menus = {
        "Student": [
            "Profile", "Materials", "XP Journey", "SOP", "Leaderboard",
            "Results", "Request XP", "Quiz",
        ],
        "Admin": [
            "User access", "CRUD", "Material", "Quiz", "Award XP", "SOP",
            "Student record", "Analysis background", "Analysis results",
            "Analysis XP",
        ],
    }
    active_key = f"active_page_{role.lower()}"
    if st.session_state.get(active_key) not in menus[role]:
        st.session_state[active_key] = current_page
    mobile_key = f"mobile_page_{role.lower()}"
    if st.session_state.get(mobile_key) not in menus[role]:
        st.session_state[mobile_key] = st.session_state[active_key]
    with st.container(key=f"mobile_navigation_{role.lower()}"):
        st.markdown('<span class="mobile-nav-marker"></span>', unsafe_allow_html=True)
        with st.container(
            horizontal=True,
            horizontal_alignment="left",
            vertical_alignment="center",
            gap="small",
        ):
            if st.button(
                "SIGN OUT",
                key=f"mobile_sign_out_{role.lower()}",
                width="content",
            ):
                st.session_state.clear()
                st.rerun()
            st.selectbox(
                "PAGE",
                menus[role],
                format_func=str.upper,
                key=mobile_key,
                on_change=sync_navigation,
                args=(mobile_key, active_key),
                label_visibility="collapsed",
                width="stretch",
            )
    return st.session_state[active_key]


def login() -> None:
    left, centre, right = st.columns([1, 1.15, 1])
    with centre:
        st.markdown("<div style='height:6vh'></div>", unsafe_allow_html=True)
        centered_image(LOGO_PATH, 132)
        client = db()
        if not client:
            st.error("XPLMS server access is not configured. Add SUPABASE_SECRET_KEY to Streamlit secrets.")
            return
        try:
            existing_users = client.table("app_users").select("id", count="exact").limit(1).execute()
            user_count = existing_users.count or 0
        except Exception:
            st.error("The app_users table is not available. Run supabase_migration_002_app_users.sql first.")
            return

        if user_count == 0:
            st.info("Create the first XPLMS administrator. This setup appears only while no application users exist.")
            with st.form("bootstrap_admin"):
                full_name = st.text_input("Administrator name")
                email = st.text_input("Administrator email")
                password = st.text_input("New password", type="password")
                confirm = st.text_input("Confirm password", type="password")
                if st.form_submit_button("Create administrator", type="primary"):
                    if not full_name.strip() or "@" not in email:
                        st.error("Enter a valid name and email.")
                    elif password != confirm:
                        st.error("The passwords do not match.")
                    elif not valid_password(password):
                        st.error("Use at least 10 characters with uppercase, lowercase and a number.")
                    else:
                        try:
                            client.table("app_users").insert({
                                "email": email.strip().lower(),
                                "username": email.strip().lower(),
                                "full_name": full_name.strip(),
                                "role": "admin",
                                "NO MATRIK": None,
                                "password_hash": hash_password(password),
                                "active": True,
                                "must_change_password": False,
                            }).execute()
                            st.success("Administrator created. Sign in below.")
                            st.rerun()
                        except Exception as exc:
                            st.error(f"Administrator could not be created: {exc}")
            return

        with st.container(border=True):
            username = st.text_input("Username")
            password = st.text_input("Password", type="password")
            if st.button("Sign in", type="primary", use_container_width=True):
                try:
                    rows = (
                        client.table("app_users")
                        .select('id,email,username,full_name,role,"NO MATRIK",password_hash,active,must_change_password')
                        .ilike("username", username.strip())
                        .limit(1)
                        .execute()
                        .data
                    )
                    account = rows[0] if rows else None
                    if (
                        not account
                        or not account.get("active")
                        or not bcrypt.checkpw(
                            password.encode("utf-8"),
                            account["password_hash"].encode("utf-8"),
                        )
                    ):
                        raise ValueError("Invalid credentials")
                    st.session_state.user = {
                        "name": account["full_name"],
                        "email": account["email"],
                        "role": "Admin" if account["role"] == "lecturer" else account["role"].title(),
                        "id": str(account["id"]),
                        "no_matrik": account.get("NO MATRIK"),
                        "must_change_password": account.get("must_change_password", False),
                    }
                    client.table("app_users").update(
                        {"last_login_at": datetime.now().isoformat()}
                    ).eq("id", account["id"]).execute()
                    st.rerun()
                except Exception:
                    st.error("We couldn't sign you in. Check your username and password.")


def merged_students() -> tuple[pd.DataFrame, bool]:
    students, live_a = fetch_table("stud_background", DEMO_STUDENTS)
    progress, live_b = fetch_table("stud_progress", DEMO_PROGRESS)
    xp, live_c = fetch_table(
        "stud_xp", DEMO_PROGRESS[["student_id", "xp"]].copy()
    )
    key = next((x for x in ["NO MATRIK", "student_id", "stud_id", "id"] if x in students.columns and x in progress.columns), None)
    if not key:
        return students, live_a
    merged = students.merge(progress, on=key, how="left")
    xp_key = next((x for x in ["NO MATRIK", "student_id", "stud_id", "id"] if x in merged.columns and x in xp.columns), None)
    if xp_key:
        merged = merged.merge(xp, on=xp_key, how="left")
    team_columns = [
        column for column in ["XPTEAM_x", "XPTEAM", "XPTEAM_y"]
        if column in merged.columns
    ]
    if team_columns:
        team = merged[team_columns[0]]
        for column in team_columns[1:]:
            team = team.combine_first(merged[column])
        merged["XPTEAM"] = team.fillna("Unassigned").astype(str)
        merged = merged.drop(
            columns=[column for column in ["XPTEAM_x", "XPTEAM_y"] if column in merged.columns]
        )
    if "XP" in merged.columns:
        merged["xp"] = pd.to_numeric(merged["XP"], errors="coerce").fillna(0)
        merged = merged.drop(columns=["XP"])
    return merged, live_a and live_b and live_c


def student_identity(user: dict[str, Any]) -> tuple[str, str]:
    """Return the student's preferred nickname and class."""
    matric = user.get("no_matrik")
    matched, _ = fetch_student_row(
        "stud_background", matric, DEMO_STUDENTS
    )
    if not matched.empty:
        row = matched.iloc[0]
        nickname = ""
        for column in ["NICKNAME PELAJAR", "nickname", "NAMA PELAJAR", "name"]:
            value = row.get(column)
            if pd.notna(value) and str(value).strip():
                nickname = str(value).strip()
                break
        student_class = row.get("KELAS", row.get("class", row.get("cohort", "—")))
        return (
            nickname or str(user.get("name", "Student")),
            "—" if pd.isna(student_class) else str(student_class),
        )
    return str(user.get("name", "Student")), "—"


ASSESSMENT_COLUMNS = [
    "C1C2",
    "C5",
    "C8",
    "C9C10",
    "INDIVIDUAL ASSIGNMENT",
    "GROUP ASSIGNMENT",
    "UPS 1",
    "UPS 2",
    "UPS 3",
    # Demo/legacy compatibility; live XPLMS uses the assessment fields above.
    "assignment",
    "pb",
    "task",
]

BADGE_LEVELS = [
    (100, "Rookie"),
    (300, "Explorer"),
    (600, "Expert"),
    (1000, "Legend"),
]
STREAK_BADGE_LEVELS = [
    (7, "Rookie"),
    (14, "Explorer"),
    (21, "Expert"),
    (28, "Legend"),
]


def streak_summary() -> tuple[pd.DataFrame, bool]:
    return fetch_table(
        "student_streak_summary",
        pd.DataFrame([
            {
                "NO MATRIK": "S24001",
                "current_streak": 8,
                "longest_streak": 8,
                "last_activity_date": "2026-07-30",
            },
            {
                "NO MATRIK": "S24002",
                "current_streak": 3,
                "longest_streak": 5,
                "last_activity_date": "2026-07-30",
            },
        ]),
    )


def badge_name_for_xp(xp: Any) -> str:
    value = float(pd.to_numeric(pd.Series([xp]), errors="coerce").fillna(0).iloc[0])
    earned = [name for threshold, name in BADGE_LEVELS if value >= threshold]
    return earned[-1] if earned else "None"


def progress_standings(
    assessment: str,
) -> tuple[pd.DataFrame, pd.DataFrame]:
    """Rank one assessment on its own; never combine unrelated marks."""
    merged, _ = merged_students()
    if merged.empty or assessment not in merged.columns:
        return pd.DataFrame(), pd.DataFrame()
    matric_col = next(
        (c for c in ["NO MATRIK", "student_id", "stud_id"] if c in merged.columns),
        None,
    )
    name_col = next(
        (c for c in ["NICKNAME PELAJAR", "nickname", "NAMA PELAJAR", "name"] if c in merged.columns),
        None,
    )
    class_col = next(
        (c for c in ["KELAS", "CLASS", "class", "cohort"] if c in merged.columns),
        None,
    )
    if not matric_col:
        return pd.DataFrame(), pd.DataFrame()
    marks = pd.to_numeric(merged[assessment], errors="coerce")
    individual = pd.DataFrame({
        "Student": merged[name_col].fillna("Unknown").astype(str) if name_col else merged[matric_col].astype(str),
        "NO MATRIK": merged[matric_col].astype(str),
        "Class": merged[class_col].fillna("Unassigned").astype(str) if class_col else "Unassigned",
        "Mark": marks,
    }).dropna(subset=["Mark"])
    individual["Rank"] = individual["Mark"].rank(method="min", ascending=False).astype(int)
    individual = individual.sort_values(["Rank", "Student"])
    classes = (
        individual.groupby("Class", dropna=False)
        .agg(
            Students=("NO MATRIK", "nunique"),
            **{
                "Average mark": ("Mark", "mean"),
                "Results available": ("Mark", "count"),
            },
        )
        .reset_index()
    )
    if not classes.empty:
        classes["Rank"] = classes["Average mark"].rank(method="min", ascending=False).astype(int)
        classes = classes.sort_values(["Rank", "Class"])
    return individual, classes


def leaderboard_data(
    month: str | None = None,
) -> tuple[pd.DataFrame, pd.DataFrame, list[str], str]:
    """Build fair individual and class rankings from the XP ledger."""
    merged, _ = merged_students()
    events, _ = fetch_table("xp_events", pd.DataFrame())
    earned_badges, badges_live = fetch_table("student_badges", pd.DataFrame())
    streaks, _ = streak_summary()
    current_month = datetime.now(ZoneInfo("Asia/Kuala_Lumpur")).strftime("%Y-%m")
    available_months = [current_month]
    if not events.empty and "created_at" in events.columns:
        event_dates = pd.to_datetime(events["created_at"], errors="coerce", utc=True)
        event_months = (
            event_dates.dt.tz_convert("Asia/Kuala_Lumpur")
            .dt.strftime("%Y-%m")
            .dropna()
            .unique()
            .tolist()
        )
        available_months = sorted(set([*event_months, current_month]), reverse=True)
    selected_month = month if month in available_months else available_months[0]
    if merged.empty:
        return pd.DataFrame(), pd.DataFrame(), available_months, selected_month

    matric_col = next(
        (column for column in ["NO MATRIK", "student_id", "stud_id"] if column in merged.columns),
        None,
    )
    name_col = next(
        (
            column
            for column in [
                "NICKNAME PELAJAR", "nickname", "NAMA PELAJAR", "name", "student_name"
            ]
            if column in merged.columns
        ),
        None,
    )
    class_col = next(
        (column for column in ["KELAS", "CLASS", "class", "cohort"] if column in merged.columns),
        None,
    )
    if not matric_col:
        return pd.DataFrame(), pd.DataFrame(), available_months, selected_month

    individuals = pd.DataFrame({
        "NO MATRIK": merged[matric_col].astype(str),
        "Student": (
            merged[name_col].fillna("Unknown").astype(str)
            if name_col
            else merged[matric_col].astype(str)
        ),
        "Class": (
            merged[class_col].fillna("Unassigned").astype(str)
            if class_col
            else "Unassigned"
        ),
        "XPTEAM": (
            merged["XPTEAM"].replace("", "Unassigned").fillna("Unassigned").astype(str)
            if "XPTEAM" in merged.columns
            else "Unassigned"
        ),
    })
    overall_values = (
        merged["xp"]
        if "xp" in merged.columns
        else pd.Series(0, index=merged.index)
    )
    individuals["Overall XP"] = pd.to_numeric(
        overall_values, errors="coerce"
    ).fillna(0)

    monthly_xp = pd.Series(dtype=float)
    if not events.empty and {"NO MATRIK", "points", "created_at"}.issubset(events.columns):
        event_dates = pd.to_datetime(events["created_at"], errors="coerce", utc=True)
        event_month = event_dates.dt.tz_convert("Asia/Kuala_Lumpur").dt.strftime("%Y-%m")
        month_events = events[event_month == selected_month].copy()
        month_events["points"] = pd.to_numeric(month_events["points"], errors="coerce").fillna(0)
        monthly_xp = month_events.groupby("NO MATRIK")["points"].sum()
    individuals["Monthly XP"] = (
        individuals["NO MATRIK"].map(monthly_xp).fillna(0)
    )

    individuals["Monthly Rank"] = (
        individuals["Monthly XP"].rank(method="min", ascending=False).astype(int)
    )
    individuals["Overall Rank"] = (
        individuals["Overall XP"].rank(method="min", ascending=False).astype(int)
    )
    if (
        badges_live
        and not earned_badges.empty
        and {"NO MATRIK", "badge_name", "xp_threshold"}.issubset(earned_badges.columns)
    ):
        xp_badges = earned_badges.copy()
        if "badge_family" in xp_badges.columns:
            xp_badges = xp_badges[xp_badges["badge_family"] == "xp"]
        permanent_badges = (
            xp_badges.sort_values("xp_threshold")
            .groupby("NO MATRIK", as_index=True)["badge_name"]
            .last()
        )
        individuals["Badge"] = (
            individuals["NO MATRIK"].map(permanent_badges).fillna("None")
        )
    else:
        individuals["Badge"] = individuals["Overall XP"].map(badge_name_for_xp)
    if not streaks.empty and {"NO MATRIK", "current_streak"}.issubset(streaks.columns):
        streak_map = streaks.set_index("NO MATRIK")["current_streak"]
        individuals["Current Streak"] = pd.to_numeric(
            individuals["NO MATRIK"].map(streak_map), errors="coerce"
        ).fillna(0).astype(int)
    else:
        individuals["Current Streak"] = 0
    if (
        badges_live
        and not earned_badges.empty
        and {"NO MATRIK", "badge_family", "badge_name", "streak_threshold"}.issubset(
            earned_badges.columns
        )
    ):
        permanent_streak_badges = (
            earned_badges[earned_badges["badge_family"] == "streak"]
            .sort_values("streak_threshold")
            .groupby("NO MATRIK", as_index=True)["badge_name"]
            .last()
        )
        individuals["Streak Badge"] = (
            individuals["NO MATRIK"].map(permanent_streak_badges).fillna("None")
        )
    else:
        individuals["Streak Badge"] = individuals["Current Streak"].map(
            lambda days: (
                [
                    name for threshold, name in STREAK_BADGE_LEVELS
                    if int(days) >= threshold
                ][-1]
                if any(
                    int(days) >= threshold
                    for threshold, _ in STREAK_BADGE_LEVELS
                )
                else "None"
            )
        )

    classes = (
        individuals.groupby("Class", dropna=False)
        .agg(
            Students=("NO MATRIK", "nunique"),
            **{
                "Monthly XP Total": ("Monthly XP", "sum"),
                "Monthly XP Average": ("Monthly XP", "mean"),
                "Overall XP Total": ("Overall XP", "sum"),
                "Overall XP Average": ("Overall XP", "mean"),
            },
        )
        .reset_index()
    )
    classes["Monthly Rank"] = (
        classes["Monthly XP Average"].rank(method="min", ascending=False).astype(int)
    )
    classes["Overall Rank"] = (
        classes["Overall XP Average"].rank(method="min", ascending=False).astype(int)
    )
    classes["Class Badge"] = classes["Overall XP Average"].map(badge_name_for_xp)
    return individuals, classes, available_months, selected_month


def progress_page() -> None:
    heading("", "Results")
    user = st.session_state.get("user", {})
    progress, _ = fetch_student_row(
        "stud_progress", user.get("no_matrik"), DEMO_PROGRESS
    )
    row = progress.iloc[0] if not progress.empty else None
    if row is None:
        st.info("No progress record is linked to this account.")
        return
    score_columns = [
        c for c in [
            "C1C2", "C5", "C8", "C9C10",
            "INDIVIDUAL ASSIGNMENT", "GROUP ASSIGNMENT",
            "UPS 1", "UPS 2", "UPS 3",
            "assignment", "pb", "task",
        ]
        if c in row.index
    ]
    results = []
    for column in score_columns:
        value = pd.to_numeric(pd.Series([row[column]]), errors="coerce").iloc[0]
        zone_column = f"{column}_ZONE"
        mark = None if pd.isna(value) else float(value)
        results.append({
            "Component": column.replace("_", " "),
            "Mark": mark,
            "Zone": row.get(zone_column, "—") if not pd.isna(row.get(zone_column, None)) else "—",
            "Status": "Available" if mark is not None else "Pending",
        })
    completed = [item["Mark"] for item in results if item["Mark"] is not None]
    if not completed:
        st.info("Marks have not been published yet.")
    cards = []
    for item in results:
        mark_text = "—" if item["Mark"] is None else f"{item['Mark']:.1f}"
        cards.append(
            '<div class="result-card">'
            f'<div class="result-title">{html.escape(str(item["Component"]).upper())}</div>'
            f'<div class="result-mark">{mark_text}</div>'
            f'<div class="result-meta">{html.escape(str(item["Zone"]).upper())} · '
            f'{html.escape(str(item["Status"]).upper())}</div></div>'
        )
    st.markdown(
        f'<div class="result-grid">{"".join(cards)}</div>',
        unsafe_allow_html=True,
    )


def _legacy_profile_page() -> None:
    heading("Student information", "My profile", "View the background information linked to your matric number.")
    background, _ = fetch_table("stud_background", DEMO_STUDENTS)
    user = st.session_state.get("user", {})
    if user.get("no_matrik") and "NO MATRIK" in background.columns:
        own = background[background["NO MATRIK"] == user["no_matrik"]]
        row = own.iloc[0] if not own.empty else None
    else:
        row = background.iloc[0] if not background.empty else None
    if row is None:
        st.info("No student background record is linked to this account.")
        return

    name = row.get(
        "NICKNAME PELAJAR",
        row.get("nickname", row.get("NAMA PELAJAR", row.get("name", user.get("name", "—")))),
    )
    matric = row.get("NO MATRIK", row.get("student_id", user.get("no_matrik", "—")))
    student_class = row.get("KELAS", row.get("cohort", "—"))
    system = row.get("SISTEM", row.get("programme", "—"))
    a, b, c, d = st.columns(4)
    with a: metric("Student name", str(name), "Registered name")
    with b: metric("No Matrik", str(matric), "Student identifier")
    with c: metric("Class", str(student_class), "Current class")
    with d: metric("System", str(system), "Academic system")

    hidden = {"id", "created_at", "updated_at"}
    details = [
        {
            "Field": str(column).replace("_", " ").title(),
            "Information": "—" if pd.isna(value) else str(value),
        }
        for column, value in row.items()
        if column not in hidden
    ]
    st.subheader("Student background")
    st.dataframe(pd.DataFrame(details), hide_index=True, width="stretch")
    st.subheader("Change password")
    st.caption("Your new password is stored as a secure hash. It cannot be viewed by Admin.")
    with st.form("student_change_password"):
        current_password = st.text_input("Current password", type="password")
        new_password = st.text_input("New password", type="password")
        confirm_password = st.text_input("Confirm new password", type="password")
        if st.form_submit_button("Change password", type="primary"):
            if is_demo():
                st.success("Password change is disabled in demo mode.")
            elif new_password != confirm_password:
                st.error("The new passwords do not match.")
            elif not valid_password(new_password):
                st.error("Use at least 10 characters with uppercase, lowercase and a number.")
            else:
                client = db()
                account = (
                    client.table("app_users")
                    .select("password_hash")
                    .eq("id", user["id"])
                    .single()
                    .execute()
                    .data
                )
                if not account or not bcrypt.checkpw(
                    current_password.encode("utf-8"),
                    account["password_hash"].encode("utf-8"),
                ):
                    st.error("The current password is incorrect.")
                else:
                    client.table("app_users").update({
                        "password_hash": hash_password(new_password),
                        "must_change_password": False,
                        "password_changed_at": datetime.now().isoformat(),
                    }).eq("id", user["id"]).execute()
                    st.success("Password changed successfully.")


def profile_page() -> None:
    heading("", "Profile")
    user = st.session_state.get("user", {})
    background, _ = fetch_student_row(
        "stud_background", user.get("no_matrik"), DEMO_STUDENTS
    )
    row = background.iloc[0] if not background.empty else None

    profile_tab, password_tab = st.tabs(["PROFILE", "CHANGE PASSWORD"])
    with profile_tab:
        if row is None:
            st.info("No student background record is linked to this account.")
        else:
            name = row.get(
                "NICKNAME PELAJAR",
                row.get("NAMA PELAJAR", row.get("name", user.get("name", "—"))),
            )
            matric = row.get("NO MATRIK", user.get("no_matrik", "—"))
            values = [
                ("Student name", name),
                ("No Matrik", matric),
                ("Class", row.get("KELAS", row.get("cohort", "—"))),
                ("System", row.get("SISTEM", row.get("programme", "—"))),
                ("PKA", row.get("PKA", "—")),
            ]
            columns = st.columns(5)
            for column, (label, value) in zip(columns, values):
                with column:
                    metric(label, "—" if pd.isna(value) else str(value).upper())
            hidden = {"id", "created_at", "updated_at"}
            details = pd.DataFrame([
                {
                    "Field": str(field).replace("_", " ").upper(),
                    "Information": "—" if pd.isna(value) else str(value).upper(),
                }
                for field, value in row.items()
                if field not in hidden
            ])
            st.dataframe(details, hide_index=True, width="stretch")

    with password_tab:
        with st.form("student_change_password_tabs"):
            current_password = st.text_input("Current password", type="password")
            new_password = st.text_input("New password", type="password")
            confirm_password = st.text_input("Confirm new password", type="password")
            if st.form_submit_button("Change password", type="primary"):
                if is_demo():
                    st.success("Password change is disabled in demo mode.")
                elif new_password != confirm_password:
                    st.error("The new passwords do not match.")
                elif not valid_password(new_password):
                    st.error("Use at least 10 characters with uppercase, lowercase and a number.")
                else:
                    client = db()
                    account = (
                        client.table("app_users").select("password_hash")
                        .eq("id", user["id"]).single().execute().data
                    )
                    if not account or not bcrypt.checkpw(
                        current_password.encode("utf-8"),
                        account["password_hash"].encode("utf-8"),
                    ):
                        st.error("The current password is incorrect.")
                    else:
                        client.table("app_users").update({
                            "password_hash": hash_password(new_password),
                            "must_change_password": False,
                            "password_changed_at": datetime.now().isoformat(),
                        }).eq("id", user["id"]).execute()
                        st.success("Password changed successfully.")


XP_REQUEST_TYPES = [
    "consultation",
    "class_participation",
    "commitment",
    "study_group",
    "extra_practice",
]
XP_DAILY_REQUEST_QUOTAS = {
    "study_group": 2,
    "extra_practice": 2,
}


def xp_requests_today(claims: pd.DataFrame, claim_type: str) -> int:
    if (
        claims.empty
        or not {"claim_type", "created_at"}.issubset(claims.columns)
    ):
        return 0
    dates = pd.to_datetime(claims["created_at"], errors="coerce", utc=True)
    malaysia_dates = dates.dt.tz_convert("Asia/Kuala_Lumpur").dt.date
    today = datetime.now(ZoneInfo("Asia/Kuala_Lumpur")).date()
    return int(
        (
            (claims["claim_type"].astype(str) == claim_type)
            & (malaysia_dates == today)
        ).sum()
    )


def _legacy_xp_badge_page() -> None:
    heading(
        "Recognition request",
        "XP",
        "Request XP for a verified learning activity.",
    )
    user = st.session_state.get("user", {})
    claims, live = fetch_table("xp_claims", pd.DataFrame())
    if live and "NO MATRIK" in claims.columns:
        claims = claims[claims["NO MATRIK"].astype(str) == str(user.get("no_matrik"))]
    history, history_live = fetch_table("xp_events", pd.DataFrame())
    if history_live and "NO MATRIK" in history.columns:
        history = history[history["NO MATRIK"].astype(str) == str(user.get("no_matrik"))]
    individuals, _, _, _ = leaderboard_data()
    own = individuals[
        individuals["NO MATRIK"].astype(str) == str(user.get("no_matrik"))
    ] if not individuals.empty else pd.DataFrame()
    current_xp = int(own.iloc[0]["Overall XP"]) if not own.empty else 0

    request_tab, list_tab, history_tab, badge_tab = st.tabs(
        ["REQUEST XP", "MY XP REQUEST", "MY XP HISTORY", "MY BADGE"]
    )
    with request_tab:
        st.write(
            "Submit evidence for consultation, class participation, commitment or "
            "a study group, or extra practice such as past-year questions. "
            "An Admin will review the request."
        )
        with st.form("xp_proof_claim", clear_on_submit=True):
            claim_type = st.selectbox(
                "XP type",
                XP_REQUEST_TYPES,
                format_func=lambda value: value.replace("_", " ").title(),
            )
            quota = XP_DAILY_REQUEST_QUOTAS.get(claim_type)
            if quota:
                used = xp_requests_today(claims, claim_type)
                st.caption(f"Daily quota: {used}/{quota} requests used today.")
            title = st.text_input("Request title")
            description = st.text_area("Describe the activity and what you learned")
            proof = st.file_uploader("Proof image (optional)", type=["jpg", "jpeg", "png", "webp"])
            if st.form_submit_button("Send for approval", type="primary"):
                if not title.strip() or not description.strip():
                    st.error("Complete the title and description.")
                elif quota and xp_requests_today(claims, claim_type) >= quota:
                    st.error(
                        f"Daily quota reached: maximum {quota} "
                        f"{claim_type.replace('_', ' ')} requests per day."
                    )
                elif is_demo():
                    st.success("Demo request submitted for Admin approval.")
                else:
                    client = db()
                    path = None
                    if proof:
                        path = f"{user['id']}/{datetime.now().timestamp()}-{proof.name}"
                        ok, result = upload_file("xp-proofs", path, proof.getvalue(), proof.type)
                        if not ok:
                            st.error(result)
                            st.stop()
                    try:
                        client.table("xp_claims").insert({
                            "student_user_id": user["id"],
                            "NO MATRIK": user["no_matrik"],
                            "claim_type": claim_type,
                            "title": title.strip(),
                            "description": description.strip(),
                            "proof_path": path,
                            "status": "pending",
                        }).execute()
                        st.success("Request submitted. An Admin will review it.")
                    except Exception as exc:
                        st.error(f"Request could not be submitted: {exc}")
    with list_tab:
        if claims.empty:
            st.info("No XP requests submitted yet.")
        else:
            visible = [
                column for column in
                ["created_at", "claim_type", "title", "status", "admin_note"]
                if column in claims.columns
            ]
            st.dataframe(
                claims[visible].sort_values("created_at", ascending=False),
                hide_index=True, width="stretch",
            )
    with history_tab:
        if history.empty:
            st.info("No XP has been recorded yet.")
        else:
            visible = [
                c for c in ["created_at", "rule_code", "points", "reason", "award_mode", "source_id"]
                if c in history.columns
            ]
            st.dataframe(
                history[visible].sort_values("created_at", ascending=False),
                hide_index=True, width="stretch",
            )
    with badge_tab:
        thresholds = BADGE_LEVELS
        captured = [(points, badge) for points, badge in thresholds if current_xp >= points]
        next_badge = next(((points, badge) for points, badge in thresholds if current_xp < points), None)
        a, b = st.columns(2)
        with a: metric("Current XP", f"{current_xp:,}", "Overall cumulative")
        with b:
            metric(
                "Current badge",
                captured[-1][1] if captured else "None",
                f"{len(captured)} captured",
            )
        if next_badge:
            remaining = next_badge[0] - current_xp
            st.subheader(f"{remaining:,} XP to {next_badge[1]}")
            previous = captured[-1][0] if captured else 0
            span = next_badge[0] - previous
            st.progress(min(1.0, max(0.0, (current_xp - previous) / span)))
        else:
            st.success("All current badges captured.")
        if captured:
            st.dataframe(
                pd.DataFrame(captured, columns=["XP threshold", "Badge"]),
                hide_index=True, width="stretch",
            )


def _legacy_student_leaderboard_page() -> None:
    heading(
        "Healthy competition",
        "Leaderboard",
        "Compare monthly and overall XP, plus academic progress.",
    )
    _, _, months, default_month = leaderboard_data()
    selected_month = st.selectbox(
        "XP month",
        months,
        index=months.index(default_month),
    )
    individuals, classes, _, _ = leaderboard_data(selected_month)
    if individuals.empty:
        st.info("Leaderboard data is not available.")
        return
    user = st.session_state.user
    own_matric = str(user.get("no_matrik", ""))
    own = individuals[individuals["NO MATRIK"].astype(str) == own_matric]
    if own.empty and is_demo():
        own = individuals.head(1)
    if not own.empty:
        row = own.iloc[0]
        badge_thresholds = BADGE_LEVELS
        captured = [
            badge for threshold, badge in badge_thresholds
            if float(row["Overall XP"]) >= threshold
        ]
        a, b, c, d = st.columns(4)
        with a: metric("Monthly XP", f"{int(row['Monthly XP']):,}", selected_month)
        with b: metric("Overall XP", f"{int(row['Overall XP']):,}", "Cumulative")
        with c: metric("Badges captured", str(len(captured)), row["Badge"])
        with d: metric("Monthly rank", f"#{int(row['Monthly Rank'])}", selected_month)
        st.subheader("Captured badges")
        if captured:
            badge_columns = st.columns(len(captured))
            for column, badge in zip(badge_columns, captured):
                with column:
                    st.markdown(
                        f'<div class="metric-card" style="text-align:center">'
                        f'<div class="metric-value">✓</div><b>{badge}</b>'
                        f'<div class="metric-note">Captured</div></div>',
                        unsafe_allow_html=True,
                    )
        else:
            st.info("No badges captured yet. The first badge unlocks at 100 overall XP.")

    xp_tab, progress_tab = st.tabs(["XP leaderboard", "Progress leaderboard"])
    with xp_tab:
        st.subheader(f"Individual · {selected_month}")
        monthly = individuals.sort_values(
            ["Monthly Rank", "Overall Rank", "Student"]
        )[
            [
                "Monthly Rank", "Student", "NO MATRIK", "Class",
                "Monthly XP", "Overall XP", "Badge",
            ]
        ]
        st.dataframe(monthly, hide_index=True, width="stretch")
        st.subheader("Individual · Overall")
        overall = individuals.sort_values(["Overall Rank", "Student"])[
            [
                "Overall Rank", "Student", "NO MATRIK", "Class",
                "Overall XP", "Monthly XP", "Badge",
            ]
        ]
        st.dataframe(overall, hide_index=True, width="stretch")
        st.subheader(f"Class · {selected_month}")
        class_monthly = classes.sort_values(["Monthly Rank", "Class"])[
            [
                "Monthly Rank", "Class", "Students", "Monthly XP Average",
                "Monthly XP Total", "Overall XP Average", "Class Badge",
            ]
        ]
        st.dataframe(class_monthly, hide_index=True, width="stretch")
        st.subheader("Class · Overall")
        class_overall = classes.sort_values(["Overall Rank", "Class"])[
            [
                "Overall Rank", "Class", "Students", "Overall XP Average",
                "Overall XP Total", "Monthly XP Average", "Class Badge",
            ]
        ]
        st.dataframe(class_overall, hide_index=True, width="stretch")
    with progress_tab:
        merged, _ = merged_students()
        assessments = [c for c in ASSESSMENT_COLUMNS if c in merged.columns]
        if not assessments:
            st.info("No assessment results are available.")
        else:
            selected_assessment = st.selectbox(
                "Assessment", assessments, key="student_progress_leaderboard_assessment"
            )
            progress, class_progress = progress_standings(selected_assessment)
            st.caption("Each assessment is ranked independently. Different assessments are never averaged together.")
            st.subheader(f"Individual · {selected_assessment}")
            st.dataframe(progress, hide_index=True, width="stretch")
            st.subheader(f"Class · {selected_assessment}")
            st.caption("Classes are ranked by the median mark for this assessment.")
            st.dataframe(class_progress, hide_index=True, width="stretch")


def request_xp_page() -> None:
    heading("", "Request XP")
    user = st.session_state.get("user", {})
    claims, live = fetch_table("xp_claims", pd.DataFrame())
    if live and "NO MATRIK" in claims.columns:
        claims = claims[
            claims["NO MATRIK"].astype(str) == str(user.get("no_matrik"))
        ]

    request_tab, list_tab = st.tabs(["REQUEST XP", "MY XP REQUEST"])
    with request_tab:
        st.info(
            "Study group and Extra practice requests are limited to 2 requests "
            "per type per day. Extra practice includes past-year questions and "
            "other additional exercises."
        )
        with st.form("student_xp_request", clear_on_submit=True):
            claim_type = st.selectbox(
                "XP type",
                XP_REQUEST_TYPES,
                format_func=lambda value: value.replace("_", " ").title(),
            )
            quota = XP_DAILY_REQUEST_QUOTAS.get(claim_type)
            if quota:
                used = xp_requests_today(claims, claim_type)
                st.caption(f"Daily quota: {used}/{quota} requests used today.")
            title = st.text_input("Request title")
            description = st.text_area("Describe the activity")
            proof = st.file_uploader(
                "Proof image (optional)", type=["jpg", "jpeg", "png", "webp"]
            )
            if st.form_submit_button("Send for approval", type="primary"):
                if not title.strip() or not description.strip():
                    st.error("Complete the title and description.")
                elif quota and xp_requests_today(claims, claim_type) >= quota:
                    st.error(
                        f"Daily quota reached: maximum {quota} "
                        f"{claim_type.replace('_', ' ')} requests per day."
                    )
                elif is_demo():
                    st.success("Demo request submitted.")
                else:
                    client = db()
                    path = None
                    if proof:
                        path = f"{user['id']}/{datetime.now().timestamp()}-{proof.name}"
                        ok, result = upload_file(
                            "xp-proofs", path, proof.getvalue(), proof.type
                        )
                        if not ok:
                            st.error(result)
                            st.stop()
                    try:
                        client.table("xp_claims").insert({
                            "student_user_id": user["id"],
                            "NO MATRIK": user["no_matrik"],
                            "claim_type": claim_type,
                            "title": title.strip(),
                            "description": description.strip(),
                            "proof_path": path,
                            "status": "pending",
                        }).execute()
                        st.success("Request submitted for Admin review.")
                    except Exception as exc:
                        st.error(f"Request could not be submitted: {exc}")
    with list_tab:
        if claims.empty:
            st.info("No XP requests submitted yet.")
        else:
            visible = [
                column for column in
                ["created_at", "claim_type", "title", "status", "admin_note"]
                if column in claims.columns
            ]
            st.dataframe(
                claims[visible].sort_values("created_at", ascending=False),
                hide_index=True, width="stretch",
            )


def render_xp_streak_sop() -> None:
    sop = pd.DataFrame([
        {
            "XP event": "Consultation",
            "Method": "Admin manual award / Student request → Admin approval",
            "Points": "20 XP default",
            "Daily request quota": "—",
        },
        {
            "XP event": "Class participation",
            "Method": "Admin manual award / Student request → Admin approval",
            "Points": "10 XP default",
            "Daily request quota": "—",
        },
        {
            "XP event": "Commitment",
            "Method": "Admin manual award / Student request → Admin approval",
            "Points": "10 XP default",
            "Daily request quota": "—",
        },
        {
            "XP event": "Study group",
            "Method": "Student request → Admin approval",
            "Points": "5 XP default",
            "Daily request quota": "2 per day",
        },
        {
            "XP event": "Extra practice",
            "Method": "Student request → Admin approval",
            "Points": "15 XP default",
            "Daily request quota": "2 per day",
        },
        {
            "XP event": "In-app quiz attempt",
            "Method": "Automatic",
            "Points": "5 XP per answered question",
            "Daily request quota": "—",
        },
        {
            "XP event": "Correct quiz answer",
            "Method": "Automatic bonus",
            "Points": "Additional 5 XP",
            "Daily request quota": "—",
        },
    ])
    st.dataframe(sop, hide_index=True, width="stretch")
    st.subheader("STREAK SOP")
    streak_sop = pd.DataFrame([
        {
            "Qualifying activity": "In-app quiz",
            "Daily requirement": "Complete at least 1 full quiz set",
            "Streak date": "Quiz completion date",
        },
        {
            "Qualifying activity": "Extra practice",
            "Daily requirement": "At least 1 request approved by Admin",
            "Streak date": "Original request date",
        },
    ])
    st.dataframe(streak_sop, hide_index=True, width="stretch")
    st.markdown(
        "- One qualifying activity is enough to capture the streak day; additional "
        "qualifying activities on the same date do not add extra days.\n"
        "- Streak days must be consecutive using the Malaysia calendar date.\n"
        "- Today remains available until the day ends. If a completed day is "
        "missed, the current streak resets to **0**.\n"
        "- Streak badges already captured remain permanent after a streak reset."
    )


def render_badge_sop() -> None:
    badges = pd.DataFrame(
        [
            {
                "Badge family": family,
                "Badge": badge,
                "Requirement": (
                    f"{threshold} overall XP"
                    if family == "XP" else f"{threshold}-day streak"
                ),
            }
            for family, levels in [
                ("XP", BADGE_LEVELS),
                ("Streak", STREAK_BADGE_LEVELS),
            ]
            for threshold, badge in levels
        ]
    )
    st.dataframe(badges, hide_index=True, width="stretch")
    st.info(
        "XP and Streak badges are captured permanently when their threshold is "
        "first reached. Later XP deductions or a streak reset do not remove an "
        "earned badge. A streak day requires either one completed in-app quiz set "
        "or one approved Extra practice request."
    )


def sop_page() -> None:
    heading("", "SOP")
    xp_tab, badge_tab = st.tabs(["XP & STREAK SOP", "BADGE SOP"])
    with xp_tab:
        render_xp_streak_sop()
    with badge_tab:
        render_badge_sop()


def my_xp_page() -> None:
    heading("", "XP Journey")
    user = st.session_state.get("user", {})
    individuals, _, months, default_month = leaderboard_data()
    selected_month = st.selectbox(
        "XP month", months, index=months.index(default_month),
    )
    if selected_month != default_month:
        individuals, _, _, _ = leaderboard_data(selected_month)
    own = individuals[
        individuals["NO MATRIK"].astype(str) == str(user.get("no_matrik"))
    ] if not individuals.empty else pd.DataFrame()
    if own.empty and is_demo():
        own = individuals.head(1)
    row = own.iloc[0] if not own.empty else None
    current_xp = int(row["Overall XP"]) if row is not None else 0
    if row is None:
        xp_record, _ = fetch_student_row(
            "stud_xp",
            user.get("no_matrik"),
            pd.DataFrame(columns=["NO MATRIK", "XP"]),
        )
        if not xp_record.empty:
            xp_value = xp_record.iloc[0].get(
                "XP", xp_record.iloc[0].get("xp", 0)
            )
            numeric_xp = pd.to_numeric(xp_value, errors="coerce")
            current_xp = int(numeric_xp) if pd.notna(numeric_xp) else 0
    history, history_live = fetch_table("xp_events", pd.DataFrame())
    if history_live and "NO MATRIK" in history.columns:
        history = history[
            history["NO MATRIK"].astype(str) == str(user.get("no_matrik"))
        ]

    earned, earned_live = fetch_table("student_badges", pd.DataFrame())
    if earned_live and "NO MATRIK" in earned.columns:
        earned = earned[
            earned["NO MATRIK"].astype(str) == str(user.get("no_matrik"))
        ]
    current_streak = int(row.get("Current Streak", 0)) if row is not None else 0
    if row is None:
        streaks, _ = streak_summary()
        if not streaks.empty and "NO MATRIK" in streaks.columns:
            own_streak = streaks[
                streaks["NO MATRIK"].astype(str)
                == str(user.get("no_matrik"))
            ]
            if not own_streak.empty:
                streak_value = pd.to_numeric(
                    own_streak.iloc[0].get("current_streak", 0),
                    errors="coerce",
                )
                current_streak = (
                    int(streak_value) if pd.notna(streak_value) else 0
                )
    if earned_live and not earned.empty:
        if "badge_family" in earned.columns:
            xp_names = set(
                earned.loc[earned["badge_family"] == "xp", "badge_name"].astype(str)
            )
            streak_names = set(
                earned.loc[
                    earned["badge_family"] == "streak", "badge_name"
                ].astype(str)
            )
        else:
            xp_names = set(earned.get("badge_name", pd.Series(dtype=str)).astype(str))
            streak_names: set[str] = set()
    else:
        xp_names = {
            badge for threshold, badge in BADGE_LEVELS
            if current_xp >= threshold
        }
        streak_names = {
            badge for threshold, badge in STREAK_BADGE_LEVELS
            if current_streak >= threshold
        }

    def badge_checklist(
        levels: list[tuple[int, str]],
        captured_names: set[str],
        unit: str,
    ) -> pd.DataFrame:
        return pd.DataFrame([
            {
                "Badge": badge,
                "Requirement": f"{threshold:,} {unit}",
                "Status": "CAPTURED" if badge in captured_names else "LOCKED",
            }
            for threshold, badge in levels
        ])

    history_tab, xp_tab, streak_tab = st.tabs(
        ["XP HISTORY", "MY XP", "MY STREAK"]
    )
    with history_tab:
        if history.empty:
            st.info("No XP has been recorded yet.")
        else:
            visible = [
                c for c in
                ["created_at", "rule_code", "points", "reason", "award_mode"]
                if c in history.columns
            ]
            st.dataframe(
                history[visible].sort_values("created_at", ascending=False),
                hide_index=True, width="stretch",
            )
    with xp_tab:
        captured_xp = [
            (threshold, badge) for threshold, badge in BADGE_LEVELS
            if badge in xp_names
        ]
        current_badge = captured_xp[-1][1] if captured_xp else "None"
        next_xp = next(
            (
                (threshold, badge) for threshold, badge in BADGE_LEVELS
                if badge not in xp_names
            ),
            None,
        )
        a, b, c = st.columns(3)
        with a: metric("Current XP", f"{current_xp:,}")
        with b: metric("Current XP badge", current_badge)
        with c:
            metric(
                "Next XP badge",
                (
                    f"{max(0, next_xp[0] - current_xp):,} XP"
                    if next_xp else "Complete"
                ),
                next_xp[1] if next_xp else "All levels captured",
            )
        st.subheader("XP badge checklist")
        st.dataframe(
            badge_checklist(BADGE_LEVELS, xp_names, "XP"),
            hide_index=True, width="stretch",
        )
    with streak_tab:
        captured_streak = [
            (threshold, badge) for threshold, badge in STREAK_BADGE_LEVELS
            if badge in streak_names
        ]
        current_streak_badge = (
            captured_streak[-1][1] if captured_streak else "None"
        )
        next_streak = next(
            (
                (threshold, badge) for threshold, badge in STREAK_BADGE_LEVELS
                if badge not in streak_names
            ),
            None,
        )
        a, b, c = st.columns(3)
        with a: metric("Current streak", f"{current_streak} days")
        with b: metric("Current streak badge", current_streak_badge)
        with c:
            metric(
                "Next streak badge",
                (
                    f"{max(0, next_streak[0] - current_streak)} days"
                    if next_streak else "Complete"
                ),
                next_streak[1] if next_streak else "All levels captured",
            )
        st.subheader("Streak badge checklist")
        st.dataframe(
            badge_checklist(STREAK_BADGE_LEVELS, streak_names, "days"),
            hide_index=True, width="stretch",
        )


def student_leaderboard_page() -> None:
    heading("", "Leaderboard")
    individuals, classes, months, default_month = leaderboard_data()
    selected_month = st.selectbox(
        "XP month", months, index=months.index(default_month),
        key="leaderboard_xp_month",
    )
    if selected_month != default_month:
        individuals, classes, _, _ = leaderboard_data(selected_month)
    class_options = (
        sorted(individuals["Class"].dropna().astype(str).unique().tolist())
        if not individuals.empty and "Class" in individuals.columns
        else []
    )
    selected_class = st.selectbox(
        "Kelas", ["ALL", *class_options], key="student_leaderboard_class"
    )
    merged, _ = merged_students()
    excluded = {
        "UPS 1", "UPS 2", "UPS 3",
        "INDIVIDUAL ASSIGNMENT", "GROUP ASSIGNMENT",
    }
    assessments = [
        c for c in ASSESSMENT_COLUMNS
        if c in merged.columns and c not in excluded
    ]
    assessment = st.selectbox(
        "Test", assessments,
        key="student_leaderboard_assessment",
    ) if assessments else None
    individual_progress, class_progress = (
        progress_standings(assessment)
        if assessment else (pd.DataFrame(), pd.DataFrame())
    )
    if selected_class != "ALL":
        individuals = individuals[
            individuals["Class"].astype(str) == selected_class
        ]
        classes = classes[classes["Class"].astype(str) == selected_class]
        individual_progress = individual_progress[
            individual_progress["Class"].astype(str) == selected_class
        ]
        class_progress = class_progress[
            class_progress["Class"].astype(str) == selected_class
        ]

    teams = (
        individuals.groupby("XPTEAM", dropna=False)
        .agg(
            Students=("Student", "count"),
            **{
                "Monthly XP Average": ("Monthly XP", "mean"),
                "Monthly XP Total": ("Monthly XP", "sum"),
                "Overall XP Average": ("Overall XP", "mean"),
                "Overall XP Total": ("Overall XP", "sum"),
            },
        )
        .reset_index()
        if not individuals.empty and "XPTEAM" in individuals.columns
        else pd.DataFrame()
    )
    if not teams.empty:
        teams["Monthly Rank"] = teams["Monthly XP Average"].rank(
            method="min", ascending=False
        ).astype(int)
        teams["Overall Rank"] = teams["Overall XP Average"].rank(
            method="min", ascending=False
        ).astype(int)

    xp_individual, xp_class, xp_team, progress_individual, progress_class = st.tabs([
        "XP & STREAK (IND)", "XP (CLASS)",
        "XP (XPTEAM)", "TEST (IND)", "TEST (CLASS)",
    ])
    with xp_individual:
        if individuals.empty:
            st.info("Individual XP leaderboard is not available.")
        else:
            st.dataframe(
                individuals[[
                    "Monthly Rank", "Overall Rank", "Student", "Class",
                    "Monthly XP", "Overall XP", "Current Streak", "Streak Badge",
                ]].sort_values(["Monthly Rank", "Overall Rank", "Student"]),
                hide_index=True, width="stretch",
            )
    with xp_class:
        if classes.empty:
            st.info("Class XP leaderboard is not available.")
        else:
            st.dataframe(
                classes[[
                    "Monthly Rank", "Overall Rank", "Class", "Students",
                    "Monthly XP Average", "Monthly XP Total",
                    "Overall XP Average", "Overall XP Total",
                ]].sort_values(["Monthly Rank", "Overall Rank", "Class"]),
                hide_index=True, width="stretch",
            )
    with xp_team:
        if teams.empty:
            st.info("XPTEAM leaderboard is not available.")
        else:
            st.dataframe(
                teams[[
                    "Monthly Rank", "Overall Rank", "XPTEAM", "Students",
                    "Monthly XP Average", "Monthly XP Total",
                    "Overall XP Average", "Overall XP Total",
                ]].sort_values(["Monthly Rank", "Overall Rank", "XPTEAM"]),
                hide_index=True, width="stretch",
            )
    with progress_individual:
        if individual_progress.empty:
            st.info("Individual progress leaderboard is not available.")
        else:
            st.dataframe(
                individual_progress[["Rank", "Student", "Class", "Mark"]],
                hide_index=True, width="stretch",
            )
    with progress_class:
        if class_progress.empty:
            st.info("Class progress leaderboard is not available.")
        else:
            st.dataframe(
                class_progress[[
                    "Rank", "Class", "Students", "Average mark",
                ]],
                hide_index=True, width="stretch",
            )


def render_admin_material_view(
    materials: pd.DataFrame,
    live: bool,
    chapters: list[int],
    material_types: list[str],
) -> None:
    filter_a, filter_b = st.columns(2)
    with filter_a:
        chapter_filter = st.segmented_control(
            "Chapter", ["All", *chapters], default="All",
            format_func=lambda value: "ALL" if value == "All" else f"C{value}",
            key="admin_view_material_chapter",
        )
    with filter_b:
        type_filter = st.selectbox(
            "Material type", ["All types", *material_types],
            key="admin_view_material_type",
        )
    query = st.text_input(
        "Search materials", placeholder="Search by title, subject or type…",
        key="admin_view_material_search",
    )
    shown = materials.copy()
    if query:
        shown = shown[
            shown.astype(str).apply(
                lambda row: row.str.contains(query, case=False).any(), axis=1
            )
        ]
    if chapter_filter != "All" and "chapter" in shown.columns:
        shown = shown[
            pd.to_numeric(shown["chapter"], errors="coerce") == int(chapter_filter)
        ]
    if type_filter != "All types" and "material_type" in shown.columns:
        shown = shown[shown["material_type"] == type_filter]
    if shown.empty:
        st.info("No materials match the selected filters.")
        return
    for index, item in shown.iterrows():
        with st.container(border=True):
            c1, c2, c3 = st.columns([5, 2.4, 1.4])
            c1.markdown(f"**{item['title']}**  \n{item['course']}")
            c2.caption(
                f"C{item.get('chapter', '—')} · "
                f"{item.get('material_type', 'Other')} · {item.get('type', 'FILE')}"
            )
            if live and item.get("file_path"):
                try:
                    signed = db().storage.from_("materials").create_signed_url(
                        str(item["file_path"]), 600
                    )
                    url = signed.get("signedURL") or signed.get("signedUrl")
                    if not url:
                        raise ValueError("No signed URL returned.")
                    c3.link_button("View", url, width="stretch")
                except Exception:
                    c3.button(
                        "Unavailable", key=f"admin_view_material_{index}",
                        disabled=True, width="stretch",
                    )


def materials_page(lecturer: bool = False, quiz_tools: bool = False) -> None:
    heading(
        "Knowledge library",
        "Study materials",
        "Browse learning resources by chapter and material type.",
    )
    materials, live = fetch_table("materials", EMPTY_MATERIALS)
    chapters = [1, 2, 5, 8, 9, 10]
    material_types = ["Infographic", "Notes", "Exercise", "Extra", "Reference", "Other"]

    if lecturer:
        view_tab, upload_tab, manage_tab = st.tabs(["VIEW", "UPLOAD", "EDIT & DELETE"])
        with view_tab:
            render_admin_material_view(
                materials, live, chapters, material_types
            )
        with upload_tab:
            with st.form("material_upload"):
                title = st.text_input("Resource title")
                course = st.text_input("Course or subject", value="Mathematics")
                chapter = st.segmented_control(
                    "Chapter",
                    chapters,
                    default=chapters[0],
                    format_func=lambda value: f"C{value}",
                )
                material_type = st.segmented_control(
                    "Material type",
                    material_types,
                    default=material_types[0],
                )
                files = st.file_uploader(
                    "Files",
                    type=[
                        "pdf", "docx", "pptx", "xlsx", "txt", "md", "csv", "zip",
                        "jpg", "jpeg", "png", "webp", "gif", "bmp", "tif", "tiff",
                    ],
                    accept_multiple_files=True,
                )
                if st.form_submit_button("Share with students", type="primary"):
                    if not title.strip():
                        st.error("Enter a resource title.")
                    elif not files:
                        st.error("Choose at least one file.")
                    else:
                        client = db()
                        try:
                            client.table("materials").select("chapter,material_type").limit(1).execute()
                        except Exception:
                            st.error("Run supabase_migration_005_material_classification.sql before uploading materials.")
                        else:
                            uploaded_paths: list[str] = []
                            rows: list[dict[str, Any]] = []
                            error_message = ""
                            for index, file in enumerate(files):
                                safe_name = re.sub(
                                    r"[^A-Za-z0-9._-]+", "-", file.name
                                ).strip("-") or f"file-{index + 1}"
                                storage_path = (
                                    f"chapter-{chapter}/{material_type.lower()}/"
                                    f"{datetime.now().timestamp()}-{index}-{safe_name}"
                                )
                                ok, path = upload_file(
                                    "materials", storage_path, file.getvalue(),
                                    file.type or "application/octet-stream",
                                )
                                if not ok:
                                    error_message = path
                                    break
                                uploaded_paths.append(path)
                                rows.append({
                                        "title": title.strip(),
                                        "course": course.strip() or "Mathematics",
                                        "chapter": int(chapter),
                                        "material_type": material_type,
                                        "file_path": path,
                                        "type": file.name.rsplit(".", 1)[-1].upper(),
                                        "lecturer_id": st.session_state.user["id"],
                                })
                            if error_message:
                                if uploaded_paths:
                                    try:
                                        client.storage.from_("materials").remove(uploaded_paths)
                                    except Exception:
                                        pass
                                st.error(error_message)
                            else:
                                saved, message = upsert_rows(
                                    "materials", pd.DataFrame(rows)
                                )
                                if saved:
                                    st.success(
                                        f"{len(rows)} file(s) shared with students."
                                    )
                                    st.rerun()
                                else:
                                    try:
                                        client.storage.from_("materials").remove(uploaded_paths)
                                    except Exception:
                                        pass
                                    st.error(message)
        with manage_tab:
            if materials.empty or "id" not in materials.columns:
                st.info("No uploaded materials are available to manage.")
            else:
                material_labels = {
                    row["id"]: (
                        f"{row.get('title', 'Untitled')} · "
                        f"{Path(str(row.get('file_path', 'file'))).name}"
                    )
                    for _, row in materials.iterrows()
                }
                selected_id = st.selectbox(
                    "Material entry",
                    list(material_labels),
                    format_func=lambda value: material_labels[value],
                    key="manage_material_id",
                )
                selected = materials[materials["id"] == selected_id].iloc[0]
                st.markdown(
                    f"**{selected.get('title', 'Untitled')}**  \n"
                    f"C{selected.get('chapter', '—')} · "
                    f"{selected.get('material_type', 'Other')} · "
                    f"{selected.get('type', 'FILE')}"
                )
                if live and selected.get("file_path"):
                    try:
                        signed = db().storage.from_("materials").create_signed_url(
                            str(selected["file_path"]), 600
                        )
                        selected_url = (
                            signed.get("signedURL") or signed.get("signedUrl")
                        )
                        if selected_url:
                            st.link_button("View selected file", selected_url)
                    except Exception:
                        st.warning("The selected file is currently unavailable.")
                with st.form("edit_material"):
                    edited_title = st.text_input(
                        "Resource title", value=str(selected.get("title", ""))
                    )
                    edited_course = st.text_input(
                        "Course or subject", value=str(selected.get("course", "Mathematics"))
                    )
                    current_chapter = pd.to_numeric(
                        pd.Series([selected.get("chapter")]), errors="coerce"
                    ).iloc[0]
                    edited_chapter = st.selectbox(
                        "Chapter", chapters,
                        index=chapters.index(int(current_chapter))
                        if pd.notna(current_chapter) and int(current_chapter) in chapters
                        else 0,
                        format_func=lambda value: f"C{value}",
                    )
                    current_type = str(selected.get("material_type", material_types[0]))
                    edited_type = st.selectbox(
                        "Material type", material_types,
                        index=material_types.index(current_type)
                        if current_type in material_types else 0,
                    )
                    if st.form_submit_button("Save changes", type="primary"):
                        if not edited_title.strip():
                            st.error("Enter a resource title.")
                        elif not live:
                            st.info("Demo mode is read-only.")
                        else:
                            try:
                                db().table("materials").update({
                                    "title": edited_title.strip(),
                                    "course": edited_course.strip() or "Mathematics",
                                    "chapter": int(edited_chapter),
                                    "material_type": edited_type,
                                }).eq("id", selected_id).execute()
                                st.success("Material details updated.")
                                st.rerun()
                            except Exception as exc:
                                st.error(f"Material could not be updated: {exc}")
                confirm_delete = st.checkbox(
                    "I understand this permanently deletes the uploaded file.",
                    key=f"confirm_material_delete_{selected_id}",
                )
                if st.button(
                    "Delete material", type="primary",
                    disabled=not confirm_delete,
                    key=f"delete_material_{selected_id}",
                ):
                    if not live:
                        st.info("Demo mode is read-only.")
                    else:
                        try:
                            file_path = str(selected.get("file_path", "")).strip()
                            db().table("materials").delete().eq(
                                "id", selected_id
                            ).execute()
                            if file_path:
                                db().storage.from_("materials").remove([file_path])
                            st.success("Material and its stored file were deleted.")
                            st.rerun()
                        except Exception as exc:
                            st.error(f"Material could not be deleted: {exc}")
        return

    filter_a, filter_b = st.columns(2)
    with filter_a:
        chapter_filter = st.segmented_control(
            "Chapter",
            ["All", *chapters],
            default="All",
            format_func=lambda value: (
                "ALL" if value == "All" else f"C{value}"
            ),
            key=f"{'admin' if lecturer else 'student'}_material_chapter",
        )
    with filter_b:
        type_filter = st.selectbox(
            "Material type",
            ["All types", *material_types],
            key=f"{'admin' if lecturer else 'student'}_material_type",
        )

    query = st.text_input("Search materials", placeholder="Search by title, subject or type…")
    if query:
        materials = materials[materials.astype(str).apply(lambda row: row.str.contains(query, case=False).any(), axis=1)]
    if chapter_filter != "All" and "chapter" in materials.columns:
        chapter_values = pd.to_numeric(materials["chapter"], errors="coerce")
        materials = materials[chapter_values == int(chapter_filter)]
    if type_filter != "All types" and "material_type" in materials.columns:
        materials = materials[materials["material_type"] == type_filter]

    if materials.empty:
        st.info(
            "No materials have been uploaded yet."
            if not query and chapter_filter == "All" and type_filter == "All types"
            else "No materials match the selected filters."
        )

    for i, item in materials.iterrows():
        with st.container(border=True):
            c1, c2, c3, c4 = st.columns([5, 2.4, 1.2, 1.6])
            c1.markdown(f"**{item['title']}**  \n{item['course']}")
            c2.caption(
                f"C{item.get('chapter', '—')} · "
                f"{item.get('material_type', 'Other')} · {item.get('type', 'FILE')}"
            )
            if live and item.get("file_path"):
                try:
                    signed = db().storage.from_("materials").create_signed_url(
                        str(item["file_path"]), 600
                    )
                    url = signed.get("signedURL") or signed.get("signedUrl")
                    if not url:
                        raise ValueError("No signed material URL returned.")
                    c3.link_button(
                        "Download",
                        url,
                        width="stretch",
                    )
                except Exception:
                    c3.button(
                        "Unavailable",
                        key=f"material_{i}",
                        disabled=True,
                        width="stretch",
                    )
            if quiz_tools and live and item.get("file_path") and c4.button("Generate quiz", key=f"quiz_material_{i}"):
                with st.spinner("Generating a draft quiz from this material…"):
                    ok, message = generate_material_quiz(item, 10)
                (st.success if ok else st.error)(message)

    if quiz_tools:
        st.subheader("Quiz management")
        quizzes, quiz_live = fetch_table("quizzes", pd.DataFrame())
        if quizzes.empty:
            st.info("No quizzes created yet.")
        else:
            for _, quiz in quizzes.iterrows():
                with st.container(border=True):
                    q1, q2, q3 = st.columns([5, 2, 1.5])
                    q1.markdown(f"**{quiz['title']}**")
                    q2.caption(f"{quiz['status'].title()} · {quiz['xp_reward']} XP")
                    if quiz["status"] == "draft" and q3.button("Publish", key=f"publish_quiz_{quiz['id']}"):
                        db().table("quizzes").update({"status": "published"}).eq("id", quiz["id"]).execute()
                        st.success("Quiz published to students.")
                        st.rerun()


def admin_quiz_page() -> None:
    heading(
        "Assessment builder",
        "AI quiz generation",
        "Generate up to 100 multiple-choice questions from an uploaded material.",
    )
    materials, live = fetch_table("materials", EMPTY_MATERIALS)
    if materials.empty:
        st.info("Upload at least one material before generating a quiz.")
    else:
        labels = {
            int(row["id"]): (
                f"C{row.get('chapter', '—')} · "
                f"{row.get('material_type', 'Other')} · {row['title']}"
            )
            for _, row in materials.iterrows()
        }
        with st.form("admin_generate_quiz"):
            material_id = st.selectbox(
                "Source material",
                list(labels),
                format_func=lambda value: labels[value],
            )
            question_count = st.number_input(
                "Number of questions",
                min_value=1,
                max_value=100,
                value=20,
                step=5,
            )
            generate = st.form_submit_button("Generate draft quiz", type="primary")
            if generate:
                selected = materials[materials["id"] == material_id].iloc[0]
                with st.spinner(
                    f"Generating {int(question_count)} questions in batches…"
                ):
                    ok, message = generate_material_quiz(
                        selected, int(question_count)
                    )
                (st.success if ok else st.error)(message)

    st.subheader("Quiz library")
    quizzes, _ = fetch_table("quizzes", pd.DataFrame())
    if quizzes.empty:
        st.info("No quizzes created yet.")
        return
    for _, quiz in quizzes.iterrows():
        with st.container(border=True):
            left, middle, right = st.columns([5, 2, 1.5])
            left.markdown(f"**{quiz['title']}**")
            middle.caption(
                f"{str(quiz['status']).title()} · {quiz.get('xp_reward', 0)} XP"
            )
            if quiz["status"] == "draft" and right.button(
                "Publish", key=f"admin_publish_quiz_{quiz['id']}"
            ):
                db().table("quizzes").update({"status": "published"}).eq(
                    "id", quiz["id"]
                ).execute()
                st.success("Quiz published to students.")
                st.rerun()


def quiz_page() -> None:
    heading(
        "Daily knowledge check",
        "Chapter quiz",
        "Answer 10 random questions per chapter each day and earn automatic XP.",
    )
    user = st.session_state.user
    today = datetime.now(ZoneInfo("Asia/Kuala_Lumpur")).date().isoformat()
    materials, materials_live = fetch_table("materials", EMPTY_MATERIALS)
    quizzes, quizzes_live = fetch_table("quizzes", pd.DataFrame())
    live = materials_live and quizzes_live

    if live:
        required_quiz_columns = {"id", "material_id", "status"}
        required_material_columns = {"id", "chapter"}
        if (
            quizzes.empty
            or materials.empty
            or not required_quiz_columns.issubset(quizzes.columns)
            or not required_material_columns.issubset(materials.columns)
        ):
            st.info("No published chapter question banks are available.")
            return
        quizzes = quizzes[quizzes["status"] == "published"].copy()
        if quizzes.empty:
            st.info("No published chapter question banks are available.")
            return
        quiz_materials = quizzes.merge(
            materials[["id", "chapter"]],
            left_on="material_id",
            right_on="id",
            how="left",
            suffixes=("_quiz", "_material"),
        )
        quiz_materials = quiz_materials.dropna(subset=["chapter"])
        chapters = sorted(quiz_materials["chapter"].astype(int).unique().tolist())
    else:
        chapters = [1]
        quiz_materials = pd.DataFrame(
            [{"id_quiz": 1, "chapter": 1, "title": "Daily Chapter 1 quiz"}]
        )

    if not chapters:
        st.info("Published quizzes must be linked to materials with a chapter.")
        return
    chapter = st.selectbox(
        "Daily chapter quiz",
        chapters,
        format_func=lambda value: f"C{value} · 10 questions today",
    )

    if live:
        client = db()
        try:
            previous = (
                client.table("quiz_attempts")
                .select("score,correct_count,total_questions,xp_awarded,completed_at")
                .eq("student_user_id", user["id"])
                .eq("chapter", int(chapter))
                .eq("attempt_date", today)
                .limit(1)
                .execute()
                .data
            )
        except Exception:
            st.error(
                "Run supabase_migration_006_revised_xp_daily_quiz.sql "
                "before using daily quizzes."
            )
            return
        if previous:
            result = previous[0]
            st.success(
                f"Today's C{chapter} quiz is complete · "
                f"{int(result['correct_count'])}/{int(result['total_questions'])} correct · "
                f"{int(result['xp_awarded'])} XP awarded."
            )
            return
        chapter_quizzes = quiz_materials[
            quiz_materials["chapter"].astype(int) == int(chapter)
        ]
        quiz_ids = chapter_quizzes["id_quiz"].astype(int).tolist()
        questions = []
        for quiz_id in quiz_ids:
            rows = (
                client.table("quiz_questions")
                .select("*")
                .eq("quiz_id", quiz_id)
                .execute()
                .data
            )
            questions.extend(rows)
    else:
        quiz_ids = [1]
        questions = [
            {
                "id": index,
                "quiz_id": 1,
                "question": f"Demo question {index}: what is {index} + 1?",
                "options": [str(index), str(index + 1), str(index + 2), str(index + 3)],
                "correct_index": 1,
                "explanation": f"{index} + 1 equals {index + 1}.",
            }
            for index in range(1, 11)
        ]

    if len(questions) < 10:
        st.info(
            f"C{chapter} currently has {len(questions)} questions. "
            "An Admin must generate at least 10 before the daily quiz opens."
        )
        return

    seed_text = f"{user['id']}:{chapter}:{today}"
    seed = int(hashlib.sha256(seed_text.encode("utf-8")).hexdigest()[:8], 16)
    daily_questions = (
        pd.DataFrame(questions).sample(n=10, random_state=seed).to_dict("records")
    )
    st.caption(
        "XP rule: 5 XP for every attempted answer, plus 5 XP for every correct answer."
    )
    with st.form(f"daily_quiz_{chapter}_{today}"):
        answers = {}
        for number, question in enumerate(daily_questions, start=1):
            options = question["options"]
            if isinstance(options, str):
                options = json.loads(options)
            answers[str(question["id"])] = st.radio(
                f"{number}. {question['question']}",
                range(len(options)),
                format_func=lambda index, opts=options: opts[index],
                index=None,
                key=f"daily_answer_{chapter}_{today}_{question['id']}",
            )
        if st.form_submit_button("Submit today's quiz", type="primary"):
            if any(answer is None for answer in answers.values()):
                st.error("Answer all 10 questions before submitting.")
            else:
                correct = sum(
                    int(answers[str(question["id"])])
                    == int(question["correct_index"])
                    for question in daily_questions
                )
                total_questions = len(daily_questions)
                score = correct / total_questions * 100
                xp_awarded = (total_questions * 5) + (correct * 5)
                if not live:
                    st.success(
                        f"Demo result: {correct}/10 correct · {xp_awarded} XP."
                    )
                else:
                    try:
                        attempt = client.table("quiz_attempts").insert({
                            "quiz_id": int(quiz_ids[0]),
                            "student_user_id": user["id"],
                            "NO MATRIK": user["no_matrik"],
                            "answers": answers,
                            "score": score,
                            "correct_count": correct,
                            "total_questions": total_questions,
                            "passed": score >= 60,
                            "chapter": int(chapter),
                            "attempt_date": today,
                            "xp_awarded": xp_awarded,
                        }).execute().data[0]
                        event = client.table("xp_events").insert({
                            "NO MATRIK": user["no_matrik"],
                            "rule_code": "quiz_completion",
                            "points": xp_awarded,
                            "source_id": f"daily-quiz-{chapter}-{today}",
                            "reason": (
                                f"Daily C{chapter} quiz: "
                                f"{correct}/{total_questions} correct"
                            ),
                            "award_mode": "automatic",
                            "awarded_by": None,
                        }).execute().data[0]
                        client.table("quiz_attempts").update({
                            "xp_event_id": event["id"]
                        }).eq("id", attempt["id"]).execute()
                        st.success(
                            f"{correct}/10 correct · {xp_awarded} XP awarded."
                        )
                        st.rerun()
                    except Exception as exc:
                        st.error(f"Quiz result could not be saved: {exc}")


def _legacy_award_xp_page() -> None:
    heading(
        "Recognition",
        "Award experience points",
        "Award consultation, class participation and commitment, or review student requests.",
    )
    students, students_live = merged_students()
    rules_fallback = pd.DataFrame(
        [
            {"code": "consultation", "name": "Consultation", "default_points": 20, "award_mode": "manual"},
            {"code": "class_participation", "name": "Class participation", "default_points": 10, "award_mode": "manual"},
            {"code": "commitment", "name": "Commitment", "default_points": 10, "award_mode": "manual"},
        ]
    )
    rules, rules_live = fetch_table("xp_rules", rules_fallback)
    if "award_mode" in rules.columns:
        rules = rules[rules["award_mode"] == "manual"]
    if "code" in rules.columns:
        rules = rules[
            rules["code"].isin(
                ["consultation", "class_participation", "commitment"]
            )
        ]
    if "active" in rules.columns:
        rules = rules[rules["active"] == True]  # noqa: E712

    matric_col = next((c for c in ["NO MATRIK", "student_id"] if c in students.columns), None)
    name_col = next((c for c in ["NAMA PELAJAR", "name", "student_name"] if c in students.columns), None)
    if not matric_col or students.empty or rules.empty:
        st.info("Student profiles and active manual XP rules are required before XP can be awarded.")
        return

    choices = students[[matric_col] + ([name_col] if name_col else [])].drop_duplicates()
    labels = {
        str(row[matric_col]): (
            f"{row[name_col]} · {row[matric_col]}" if name_col else str(row[matric_col])
        )
        for _, row in choices.iterrows()
    }
    rule_lookup = {str(row["code"]): row for _, row in rules.iterrows()}

    with st.form("manual_xp_award", clear_on_submit=True):
        matric = st.selectbox("Student", list(labels), format_func=lambda value: labels[value])
        rule_code = st.selectbox(
            "Award category",
            list(rule_lookup),
            format_func=lambda value: str(rule_lookup[value].get("name", value)),
        )
        default_points = int(rule_lookup[rule_code].get("default_points", 10))
        points = st.number_input("XP points", min_value=1, max_value=1000, value=default_points)
        reason = st.text_area(
            "Reason",
            placeholder="Give a concise, specific reason for this recognition.",
        )
        confirmed = st.checkbox("I confirm this award is accurate and appropriate.")
        submitted = st.form_submit_button(
            "Award XP", type="primary", disabled=not confirmed
        )
        if submitted:
            if not reason.strip():
                st.error("Add a reason before awarding XP.")
            elif is_demo():
                st.success(f"Demo award recorded: +{points} XP for {labels[matric]}.")
            else:
                client = db()
                try:
                    if st.session_state.user.get("role") != "Admin":
                        raise PermissionError("Only administrators can award XP.")
                    client.table("xp_events").insert({
                        "NO MATRIK": matric,
                        "rule_code": rule_code,
                        "points": int(points),
                        "source_id": f"manual-{datetime.now().timestamp()}",
                        "reason": reason.strip(),
                        "award_mode": "manual",
                        "awarded_by": st.session_state.user["id"],
                    }).execute()
                    st.success(f"{points} XP awarded to {labels[matric]}.")
                except Exception as exc:
                    st.error(f"XP could not be awarded: {exc}")

    st.subheader("Student XP requests")
    claims_fallback = pd.DataFrame(
        [{
            "id": 1, "NO MATRIK": "S24001", "claim_type": "study_group",
            "title": "Algebra study group",
            "description": "Discussed simultaneous equations with classmates.",
            "proof_path": "", "status": "pending", "created_at": "2026-07-29",
        }]
    )
    claims, claims_live = fetch_table("xp_claims", claims_fallback)
    if "status" in claims.columns:
        claims = claims[claims["status"] == "pending"]
    if claims.empty:
        st.info("No pending XP claims.")
    else:
        client = db()
        for _, claim in claims.iterrows():
            claim_label = str(claim.get("claim_type", "study_group")).replace("_", " ").title()
            with st.expander(
                f"{claim_label} · {claim.get('title', 'XP request')} · "
                f"{claim.get('NO MATRIK', '')}"
            ):
                st.write(claim.get("description", ""))
                st.caption(f"Submitted {claim.get('created_at', 'recently')}")
                proof_path = claim.get("proof_path")
                if claims_live and proof_path:
                    try:
                        signed = client.storage.from_("xp-proofs").create_signed_url(str(proof_path), 600)
                        url = signed.get("signedURL") or signed.get("signedUrl")
                        if url:
                            st.image(url, caption="Submitted proof", width=420)
                    except Exception:
                        st.warning("The proof image could not be displayed.")
                admin_note = st.text_input("Admin note", key=f"claim_note_{claim['id']}")
                approve_col, reject_col = st.columns(2)
                if approve_col.button("Approve and award XP", type="primary", key=f"approve_claim_{claim['id']}"):
                    if not claims_live:
                        st.success("Demo claim approved.")
                    else:
                        try:
                            rule_code = str(claim.get("claim_type", "study_group"))
                            rule = client.table("xp_rules").select("default_points").eq("code", rule_code).limit(1).execute().data
                            points = (
                                int(rule[0]["default_points"])
                                if rule else {
                                    "study_group": 5,
                                    "extra_practice": 15,
                                }.get(rule_code, 20)
                            )
                            event = client.table("xp_events").insert({
                                "NO MATRIK": claim["NO MATRIK"],
                                "rule_code": rule_code,
                                "points": points,
                                "source_id": f"claim-{claim['id']}",
                                "reason": claim["title"],
                                "award_mode": "manual",
                                "awarded_by": st.session_state.user["id"],
                            }).execute().data[0]
                            client.table("xp_claims").update({
                                "status": "approved",
                                "admin_note": admin_note,
                                "reviewed_by": st.session_state.user["id"],
                                "reviewed_at": datetime.now().isoformat(),
                                "xp_event_id": event["id"],
                            }).eq("id", claim["id"]).execute()
                            st.success(f"Claim approved and {points} XP awarded.")
                            st.rerun()
                        except Exception as exc:
                            st.error(f"Claim approval failed: {exc}")
                if reject_col.button("Reject claim", key=f"reject_claim_{claim['id']}"):
                    if not claims_live:
                        st.info("Demo claim rejected.")
                    else:
                        client.table("xp_claims").update({
                            "status": "rejected",
                            "admin_note": admin_note,
                            "reviewed_by": st.session_state.user["id"],
                            "reviewed_at": datetime.now().isoformat(),
                        }).eq("id", claim["id"]).execute()
                        st.info("Claim rejected.")
                        st.rerun()

    st.subheader("How XP is awarded")
    st.markdown(
        "- **Admin awards:** consultation, class participation and commitment.\n"
        "- **Admin-approved requests:** consultation, class participation, commitment, study group or extra practice.\n"
        "- **Automatic awards:** daily in-app chapter quizzes."
    )
    if not (students_live and rules_live):
        st.caption("Showing the XP award workflow with demo data.")


def award_xp_page() -> None:
    heading("", "Award XP")
    students, students_live = merged_students()
    rules_fallback = pd.DataFrame([
        {"code": "consultation", "name": "Consultation", "default_points": 20},
        {"code": "class_participation", "name": "Class participation", "default_points": 10},
        {"code": "commitment", "name": "Commitment", "default_points": 10},
    ])
    rules, rules_live = fetch_table("xp_rules", rules_fallback)
    if "award_mode" in rules.columns:
        rules = rules[rules["award_mode"] == "manual"]
    if "code" in rules.columns:
        rules = rules[rules["code"].isin(["consultation", "class_participation", "commitment"])]
    if "active" in rules.columns:
        rules = rules[rules["active"] == True]  # noqa: E712

    events_fallback = pd.DataFrame([{
        "id": 1, "NO MATRIK": "S24001", "rule_code": "consultation",
        "points": 20, "reason": "Academic consultation",
        "award_mode": "manual", "created_at": "2026-07-29T10:00:00+08:00",
    }])
    events, events_live = fetch_table("xp_events", events_fallback)

    view_tab, manual_tab, approval_tab, records_tab = st.tabs(
        [
            "VIEW", "MANUAL AWARD", "APPROVE REQUEST", "EDIT & DELETE",
        ]
    )

    matric_col = next((c for c in ["NO MATRIK", "student_id"] if c in students.columns), None)
    name_col = next(
        (c for c in ["NICKNAME PELAJAR", "NAMA PELAJAR", "name", "student_name"] if c in students.columns),
        None,
    )
    class_col = next(
        (c for c in ["KELAS", "CLASS", "class", "cohort"] if c in students.columns),
        None,
    )
    identity_columns = [
        column for column in [matric_col, name_col, class_col] if column
    ]
    choices = (
        students[identity_columns].drop_duplicates(subset=[matric_col])
        if matric_col and not students.empty else pd.DataFrame()
    )

    def identity_label(matric: Any) -> str:
        matric_text = str(matric)
        matched = (
            choices[choices[matric_col].astype(str) == matric_text]
            if matric_col and not choices.empty else pd.DataFrame()
        )
        if matched.empty:
            return f"UNKNOWN · UNASSIGNED · {matric_text}"
        row = matched.iloc[0]
        student_name = (
            str(row.get(name_col)).strip()
            if name_col and pd.notna(row.get(name_col)) else "UNKNOWN"
        )
        student_class = (
            str(row.get(class_col)).strip()
            if class_col and pd.notna(row.get(class_col)) else "UNASSIGNED"
        )
        return f"{student_name} · {student_class} · {matric_text}"

    labels = {
        str(row[matric_col]): identity_label(row[matric_col])
        for _, row in choices.iterrows()
    }
    rule_lookup = {
        str(row["code"]): row for _, row in rules.iterrows()
    } if not rules.empty else {}

    with view_tab:
        if events.empty:
            st.info("No XP records are available.")
        else:
            shown_events = events.copy()
            shown_events["Student"] = shown_events["NO MATRIK"].map(
                identity_label
            )
            shown_events["Type"] = shown_events.get(
                "rule_code", pd.Series(index=shown_events.index, dtype=str)
            ).fillna("Unknown").astype(str).str.replace(
                "_", " ", regex=False
            ).str.title()
            shown_events["Time"] = pd.to_datetime(
                shown_events.get(
                    "created_at",
                    pd.Series(index=shown_events.index, dtype="datetime64[ns]"),
                ),
                errors="coerce", utc=True,
            ).dt.tz_convert("Asia/Kuala_Lumpur").dt.strftime(
                "%Y-%m-%d %H:%M"
            )
            shown_events["Method"] = shown_events.get(
                "award_mode", pd.Series(index=shown_events.index, dtype=str)
            ).fillna("Unknown").astype(str).str.title()
            visible_columns = [
                column for column in [
                    "Time", "Student", "Type", "points", "Method", "reason",
                ] if column in shown_events.columns
            ]
            st.dataframe(
                shown_events[visible_columns].rename(columns={
                    "points": "XP",
                    "reason": "Related / reason",
                }).sort_values("Time", ascending=False),
                hide_index=True,
                width="stretch",
            )

    with manual_tab:
        if not labels or not rule_lookup:
            st.info("Student profiles and active manual XP rules are required.")
        else:
            with st.form("manual_xp_award_tabs", clear_on_submit=True):
                matric = st.selectbox("Student", list(labels), format_func=lambda value: labels[value])
                rule_code = st.selectbox(
                    "Award category", list(rule_lookup),
                    format_func=lambda value: str(rule_lookup[value].get("name", value)),
                )
                default_points = int(rule_lookup[rule_code].get("default_points", 10))
                points = st.number_input(
                    "XP points", min_value=-1000, max_value=1000,
                    value=default_points, step=1,
                )
                reason = st.text_area("Reason")
                confirmed = st.checkbox("I confirm this XP entry is accurate.")
                if st.form_submit_button("Record XP", type="primary", disabled=not confirmed):
                    if int(points) == 0:
                        st.error("XP points cannot be zero.")
                    elif not reason.strip():
                        st.error("Add a reason for this XP entry.")
                    elif is_demo():
                        st.success(f"Demo XP entry recorded: {int(points):+d} XP.")
                    else:
                        try:
                            db().table("xp_events").insert({
                                "NO MATRIK": matric,
                                "rule_code": rule_code,
                                "points": int(points),
                                "source_id": f"manual-{datetime.now().timestamp()}",
                                "reason": reason.strip(),
                                "award_mode": "manual",
                                "awarded_by": st.session_state.user["id"],
                            }).execute()
                            st.success(f"{int(points):+d} XP recorded for {labels[matric]}.")
                        except Exception as exc:
                            st.error(f"XP could not be recorded: {exc}")

    with approval_tab:
        claims_fallback = pd.DataFrame([{
            "id": 1, "NO MATRIK": "S24001", "claim_type": "study_group",
            "title": "Algebra study group", "description": "Discussed equations.",
            "proof_path": None, "status": "pending", "created_at": "2026-07-29",
        }])
        claims, claims_live = fetch_table("xp_claims", claims_fallback)
        if "status" in claims.columns:
            claims = claims[claims["status"] == "pending"]
        if claims.empty:
            st.info("No pending XP requests.")
        else:
            client = db()
            for _, claim in claims.iterrows():
                label = str(claim.get("claim_type", "study_group")).replace("_", " ").title()
                with st.expander(
                    f"{label} · {claim.get('title', 'XP request')} · "
                    f"{identity_label(claim.get('NO MATRIK', ''))}"
                ):
                    st.write(claim.get("description", ""))
                    proof_path = claim.get("proof_path")
                    if claims_live and proof_path:
                        try:
                            signed = client.storage.from_("xp-proofs").create_signed_url(str(proof_path), 600)
                            url = signed.get("signedURL") or signed.get("signedUrl")
                            if url:
                                st.image(url, caption="Submitted proof", width=420)
                        except Exception:
                            st.warning("The optional proof image could not be displayed.")
                    note = st.text_input("Admin note", key=f"tab_claim_note_{claim['id']}")
                    approve, reject = st.columns(2)
                    if approve.button("Approve", type="primary", key=f"tab_approve_{claim['id']}"):
                        if not claims_live:
                            st.success("Demo request approved.")
                        else:
                            try:
                                rule_code = str(claim.get("claim_type", "study_group"))
                                rule = client.table("xp_rules").select("default_points").eq("code", rule_code).limit(1).execute().data
                                points = (
                                    int(rule[0]["default_points"])
                                    if rule else {
                                        "study_group": 5,
                                        "extra_practice": 15,
                                    }.get(rule_code, 20)
                                )
                                event = client.table("xp_events").insert({
                                    "NO MATRIK": claim["NO MATRIK"], "rule_code": rule_code,
                                    "points": points, "source_id": f"claim-{claim['id']}",
                                    "reason": claim["title"], "award_mode": "manual",
                                    "awarded_by": st.session_state.user["id"],
                                }).execute().data[0]
                                client.table("xp_claims").update({
                                    "status": "approved", "admin_note": note,
                                    "reviewed_by": st.session_state.user["id"],
                                    "reviewed_at": datetime.now().isoformat(),
                                    "xp_event_id": event["id"],
                                }).eq("id", claim["id"]).execute()
                                st.success(f"Request approved and {points} XP awarded.")
                                st.rerun()
                            except Exception as exc:
                                st.error(f"Request approval failed: {exc}")
                    if reject.button("Reject", key=f"tab_reject_{claim['id']}"):
                        if not claims_live:
                            st.info("Demo request rejected.")
                        else:
                            client.table("xp_claims").update({
                                "status": "rejected", "admin_note": note,
                                "reviewed_by": st.session_state.user["id"],
                                "reviewed_at": datetime.now().isoformat(),
                            }).eq("id", claim["id"]).execute()
                            st.info("Request rejected.")
                            st.rerun()
    with records_tab:
        if events.empty or "id" not in events.columns:
            st.info("No XP records are available.")
        else:
            def event_points(value: Any) -> int:
                numeric = pd.to_numeric(value, errors="coerce")
                return int(numeric) if pd.notna(numeric) else 0

            event_labels = {
                row["id"]: (
                    f"{identity_label(row.get('NO MATRIK', ''))} · "
                    f"{str(row.get('rule_code', '')).replace('_', ' ').title()} · "
                    f"{event_points(row.get('points')):+d} XP"
                )
                for _, row in events.sort_values(
                    "created_at", ascending=False, na_position="last"
                ).iterrows()
            }
            event_id = st.selectbox(
                "XP record", list(event_labels),
                format_func=lambda value: event_labels[value],
                key="manage_xp_event_id",
            )
            event = events[events["id"] == event_id].iloc[0]
            with st.form("edit_xp_event"):
                edited_points = st.number_input(
                    "XP points", min_value=-1000, max_value=1000,
                    value=event_points(event.get("points")),
                    step=1,
                )
                edited_reason = st.text_area(
                    "Reason", value=str(event.get("reason", ""))
                )
                if st.form_submit_button("Save changes", type="primary"):
                    if int(edited_points) == 0:
                        st.error("XP points cannot be zero.")
                    elif not edited_reason.strip():
                        st.error("Add a reason for this XP entry.")
                    elif not events_live:
                        st.info("Demo mode is read-only.")
                    else:
                        try:
                            db().table("xp_events").update({
                                "points": int(edited_points),
                                "reason": edited_reason.strip(),
                            }).eq("id", event_id).execute()
                            st.success("XP record and student balance updated.")
                            st.rerun()
                        except Exception as exc:
                            st.error(
                                "XP record could not be updated. Run "
                                "supabase_migration_011_editable_xp_xpteam.sql first. "
                                f"Details: {exc}"
                            )
            confirm_delete = st.checkbox(
                "I understand this reverses the XP from the student's balance.",
                key=f"confirm_xp_delete_{event_id}",
            )
            if st.button(
                "Delete XP record", type="primary",
                disabled=not confirm_delete,
                key=f"delete_xp_event_{event_id}",
            ):
                if not events_live:
                    st.info("Demo mode is read-only.")
                else:
                    try:
                        db().table("xp_events").delete().eq("id", event_id).execute()
                        st.success("XP record deleted and its points reversed.")
                        st.rerun()
                    except Exception as exc:
                        st.error(
                            "XP record could not be deleted. Run "
                            "supabase_migration_011_editable_xp_xpteam.sql first. "
                            f"Details: {exc}"
                        )
    if not (students_live and rules_live):
        st.caption("Showing the XP workflow with demo data.")


def students_page(admin: bool = False) -> None:
    merged, live = merged_students()
    heading("Records", "Student records" if admin else "Students", "Search, review and maintain trusted learner information.")
    query = st.text_input("Search records", placeholder="Name, ID, email or programme…")
    shown = merged
    if query:
        shown = merged[merged.astype(str).apply(lambda row: row.str.contains(query, case=False).any(), axis=1)]
    st.dataframe(shown, use_container_width=True, hide_index=True)
    if admin:
        st.subheader("Edit records")
        edited = st.data_editor(merged, use_container_width=True, num_rows="dynamic", hide_index=True, key="student_editor")
        if st.button("Save changes", type="primary"):
            ok, message = upsert_rows("stud_background", edited[[c for c in edited.columns if c in DEMO_STUDENTS.columns]])
            (st.success if ok else st.warning)(message)
    if not live: st.caption("Sample data is shown because the live tables are unavailable.")


def admin_student_records_page() -> None:
    heading(
        "Student information",
        "Student record",
        "Combined background, progress and XP records with field-level filters.",
    )
    merged, live = merged_students()
    technical_columns = [
        column
        for column in merged.columns
        if (
            str(column).lower() in {"id", "created_at", "updated_at"}
            or str(column).lower().startswith(("id_", "created_at_", "updated_at_"))
        )
    ]
    merged = merged.drop(columns=technical_columns, errors="ignore")
    if merged.empty:
        st.info("No student records are available.")
        return
    search = st.text_input("Search records", placeholder="Name, No Matrik or any value…")
    categorical = [
        column
        for column in merged.columns
        if merged[column].nunique(dropna=True) <= 30
        and column not in {"id", "created_at", "updated_at"}
    ]
    slicer_fields = st.multiselect(
        "Dropdown slicers",
        categorical,
        max_selections=4,
        placeholder="Select fields such as class, programme, system or zone",
    )
    filtered = merged.copy()
    slicer_columns = st.columns(max(1, min(4, len(slicer_fields))))
    for index, field in enumerate(slicer_fields):
        values = sorted(filtered[field].dropna().astype(str).unique().tolist())
        selected = slicer_columns[index % len(slicer_columns)].selectbox(
            field, ["All", *values], key=f"record_slicer_{field}"
        )
        if selected != "All":
            filtered = filtered[filtered[field].astype(str) == selected]
    if search:
        filtered = filtered[
            filtered.astype(str).apply(
                lambda row: row.str.contains(search, case=False, na=False).any(),
                axis=1,
            )
        ]
    a, b = st.columns(2)
    with a:
        metric("Students shown", f"{len(filtered):,}", f"of {len(merged):,} records")
    with b:
        metric("Available fields", f"{len(merged.columns):,}", "Background, progress and XP")
    st.dataframe(filtered, hide_index=True, width="stretch")
    if not live:
        st.caption("Sample data is shown because live tables are unavailable.")


def analysis_background_page() -> None:
    heading(
        "Student intelligence",
        "Analysis background",
        "Understand the composition of the student cohort.",
    )
    background, _ = fetch_table("stud_background", DEMO_STUDENTS)
    if background.empty:
        st.info("No student background data is available.")
        return
    fields = [
        column
        for column in background.columns
        if background[column].nunique(dropna=True) <= 40
        and column not in {"id", "NO MATRIK", "created_at", "updated_at"}
    ]
    if not fields:
        st.dataframe(background, hide_index=True, width="stretch")
        return
    field = st.selectbox("Analyse field", fields)
    summary = (
        background[field]
        .fillna("Not specified")
        .astype(str)
        .value_counts()
        .rename_axis(field)
        .reset_index(name="Students")
    )
    a, b = st.columns(2)
    with a: metric("Total students", f"{len(background):,}", "Background records")
    with b: metric("Groups", f"{len(summary):,}", f"Distinct {field} values")
    st.dataframe(summary, hide_index=True, width="stretch")


def analysis_progress_page() -> None:
    heading(
        "Academic intelligence",
        "Analysis results",
        "Review assessment marks and performance zones.",
    )
    progress, _ = fetch_table("stud_progress", DEMO_PROGRESS)
    assessments = [
        column
        for column in [
            "C1C2", "C5", "C8", "C9C10", "INDIVIDUAL ASSIGNMENT",
            "GROUP ASSIGNMENT", "UPS 1", "UPS 2", "UPS 3",
        ]
        if column in progress.columns
    ]
    if not assessments:
        assessments = [
            column for column in ["assignment", "pb", "task"]
            if column in progress.columns
        ]
    if progress.empty or not assessments:
        st.info("No assessment results are available.")
        return
    assessment = st.selectbox("Assessment", assessments)
    marks = pd.to_numeric(progress[assessment], errors="coerce")
    available = marks.dropna()
    a, b, c = st.columns(3)
    with a: metric("Results available", str(len(available)), f"of {len(progress)} students")
    with b: metric("Average mark", f"{available.mean():.1f}" if not available.empty else "—", assessment)
    with c: metric("Highest mark", f"{available.max():.1f}" if not available.empty else "—", assessment)
    zone_column = f"{assessment}_ZONE"
    columns = [c for c in ["NO MATRIK", assessment, zone_column] if c in progress.columns]
    st.dataframe(
        progress[columns].sort_values(assessment, ascending=False, na_position="last"),
        hide_index=True,
        width="stretch",
    )
    if zone_column in progress.columns:
        st.subheader("Zone distribution")
        zones = (
            progress[zone_column].fillna("Not assigned").astype(str)
            .value_counts().rename_axis("Zone").reset_index(name="Students")
        )
        st.dataframe(zones, hide_index=True, width="stretch")
    individual_progress, class_progress = progress_standings(assessment)
    st.subheader(f"Individual leaderboard · {assessment}")
    st.dataframe(individual_progress, hide_index=True, width="stretch")
    st.subheader(f"Class leaderboard · {assessment}")
    st.caption("Classes are ranked by the average mark for this assessment.")
    st.dataframe(class_progress, hide_index=True, width="stretch")


def analysis_xp_page() -> None:
    heading(
        "Recognition intelligence",
        "Analysis XP",
        "Review XP balances, badges and student rank.",
    )
    _, _, months, default_month = leaderboard_data()
    selected_month = st.selectbox(
        "Reward month",
        months,
        index=months.index(default_month),
    )
    individuals, classes, _, _ = leaderboard_data(selected_month)
    if individuals.empty:
        st.info("No XP records are available.")
        return
    a, b, c, d = st.columns(4)
    with a: metric("Monthly XP", f"{int(individuals['Monthly XP'].sum()):,}", selected_month)
    with b: metric("Overall XP", f"{int(individuals['Overall XP'].sum()):,}", "Cumulative balance")
    with c: metric("Ranked students", str(len(individuals)), "Eligible individuals")
    with d:
        metric(
            "Active streaks",
            str(int((individuals["Current Streak"] > 0).sum())),
            "Students currently active",
        )

    monthly_tab, overall_tab, class_tab = st.tabs(
        ["Monthly individual", "Overall individual", "Class leaderboard"]
    )
    with monthly_tab:
        monthly = individuals.sort_values(
            ["Monthly Rank", "Overall Rank", "Student"]
        )[
            [
                "Monthly Rank", "Student", "NO MATRIK", "Class",
                "Monthly XP", "Overall XP", "Badge",
                "Current Streak", "Streak Badge",
            ]
        ]
        st.dataframe(monthly, hide_index=True, width="stretch")
    with overall_tab:
        overall = individuals.sort_values(["Overall Rank", "Student"])[
            [
                "Overall Rank", "Student", "NO MATRIK", "Class",
                "Overall XP", "Monthly XP", "Badge",
                "Current Streak", "Streak Badge",
            ]
        ]
        st.dataframe(overall, hide_index=True, width="stretch")
    with class_tab:
        st.subheader(f"Monthly class ranking · {selected_month}")
        monthly_classes = classes.sort_values(["Monthly Rank", "Class"])[
            [
                "Monthly Rank", "Class", "Students", "Monthly XP Average",
                "Monthly XP Total", "Overall XP Average", "Class Badge",
            ]
        ]
        st.dataframe(monthly_classes, hide_index=True, width="stretch")
        st.subheader("Overall class ranking")
        overall_classes = classes.sort_values(["Overall Rank", "Class"])[
            [
                "Overall Rank", "Class", "Students", "Overall XP Average",
                "Overall XP Total", "Monthly XP Average", "Class Badge",
            ]
        ]
        st.dataframe(overall_classes, hide_index=True, width="stretch")


def admin_crud_page() -> None:
    heading(
        "Data management",
        "CRUD",
        "Add, update, delete or overwrite student datasets.",
    )
    target = st.selectbox(
        "Dataset",
        ["stud_background", "stud_progress", "stud_xp"],
    )
    data, live = fetch_table(target, pd.DataFrame())
    add_tab, update_tab, delete_tab, bulk_tab = st.tabs(
        ["Add", "Update", "Delete", "Bulk overwrite"]
    )
    with add_tab:
        st.caption("Add one record using the dataset columns below.")
        template_columns = [
            column for column in data.columns
            if column not in {"id", "created_at", "updated_at"}
        ]
        if not template_columns:
            st.info("Run the related Supabase schema before adding records.")
        else:
            with st.form(f"crud_add_{target}"):
                values = {
                    column: st.text_input(column)
                    for column in template_columns
                }
                if st.form_submit_button("Add record", type="primary"):
                    record = {key: value for key, value in values.items() if value != ""}
                    try:
                        db().table(target).insert(record).execute()
                        st.success("Record added.")
                        st.rerun()
                    except Exception as exc:
                        st.error(f"Record could not be added: {exc}")
    with update_tab:
        if data.empty:
            st.info("No records are available to update.")
        else:
            base_columns = data.columns.tolist()
            update_view = data.copy()
            identification_columns: list[str] = []
            if target in {"stud_progress", "stud_xp"} and "NO MATRIK" in data.columns:
                background, _ = fetch_table("stud_background", DEMO_STUDENTS)
                background_key = next(
                    (
                        column for column in ["NO MATRIK", "student_id"]
                        if column in background.columns
                    ),
                    None,
                )
                background_name = next(
                    (
                        column for column in ["NAMA PELAJAR", "name"]
                        if column in background.columns
                    ),
                    None,
                )
                background_class = next(
                    (
                        column for column in ["KELAS", "class", "cohort"]
                        if column in background.columns
                    ),
                    None,
                )
                if background_key:
                    identity = background[[
                        column for column in [
                            background_key, background_name, background_class
                        ] if column
                    ]].drop_duplicates(subset=[background_key])
                    identity = identity.rename(columns={
                        background_key: "NO MATRIK",
                        **(
                            {background_name: "NAMA PELAJAR"}
                            if background_name else {}
                        ),
                        **(
                            {background_class: "KELAS"}
                            if background_class else {}
                        ),
                    })
                    update_view = update_view.merge(
                        identity, on="NO MATRIK", how="left"
                    )
                    identification_columns = [
                        column for column in ["NAMA PELAJAR", "KELAS"]
                        if column in update_view.columns
                        and column not in base_columns
                    ]
                    ordered_columns = [
                        "NO MATRIK", *identification_columns,
                        *[
                            column for column in update_view.columns
                            if column not in {"NO MATRIK", *identification_columns}
                        ],
                    ]
                    update_view = update_view[ordered_columns]
            edited = st.data_editor(
                update_view,
                hide_index=True,
                width="stretch",
                disabled=identification_columns,
                column_config={
                    "NAMA PELAJAR": st.column_config.TextColumn(
                        "NAMA PELAJAR", help="Read-only identification"
                    ),
                    "KELAS": st.column_config.TextColumn(
                        "KELAS", help="Read-only identification"
                    ),
                },
            )
            if st.button("Save updates", type="primary"):
                ok, message = upsert_rows(target, edited[base_columns])
                (st.success if ok else st.error)(message)
    with delete_tab:
        if data.empty:
            st.info("No records are available to delete.")
        else:
            key = "NO MATRIK" if "NO MATRIK" in data.columns else "id"
            choices = data[key].dropna().astype(str).tolist()
            selected = st.selectbox("Record to delete", choices)
            confirmed = st.checkbox("I understand this permanently deletes the selected record.")
            if st.button("Delete record", disabled=not confirmed):
                try:
                    db().table(target).delete().eq(key, selected).execute()
                    st.success("Record deleted.")
                    st.rerun()
                except Exception as exc:
                    st.error(f"Record could not be deleted: {exc}")
    with bulk_tab:
        bulk_mode = st.radio(
            "Bulk operation",
            ["Overwrite selected columns", "Replace entire dataset"],
            horizontal=True,
        )
        if bulk_mode == "Overwrite selected columns":
            st.info(
                "Rows are matched by NO MATRIK. Only the columns you select "
                "will be changed; all other database values remain unchanged."
            )
        else:
            st.warning("Replace entire dataset deletes and recreates every record.")
        fallback_columns = {
            "stud_background": [
                "NO MATRIK", "NAMA PELAJAR", "NICKNAME PELAJAR",
                "SISTEM", "SPM_MATH", "SPM_ADDMATH", "DM015", "DM025",
                "PKA", "JANTINA", "KELAS", "XPTEAM",
            ],
            "stud_progress": [
                "NO MATRIK", "C1C2", "C1C2_ZONE", "C5", "C5_ZONE",
                "C8", "C8_ZONE", "C9C10", "C9C10_ZONE",
                "INDIVIDUAL ASSIGNMENT", "GROUP ASSIGNMENT",
                "UPS 1", "UPS 2", "UPS 3",
            ],
            "stud_xp": ["NO MATRIK", "XP", "XPTEAM"],
        }
        technical_columns = {"id", "created_at", "updated_at"}
        template_columns = [
            column for column in data.columns
            if column not in technical_columns
        ] or fallback_columns[target]
        template_csv = pd.DataFrame(columns=template_columns).to_csv(
            index=False
        ).encode("utf-8-sig")
        st.download_button(
            f"Download {target} CSV template",
            data=template_csv,
            file_name=f"{target}_bulk_template.csv",
            mime="text/csv",
            key=f"crud_template_{target}",
        )
        uploaded = st.file_uploader(
            "CSV or Excel file", type=["csv", "xlsx", "xls"], key="crud_bulk_file"
        )
        if uploaded:
            frame = (
                pd.read_csv(uploaded)
                if uploaded.name.lower().endswith(".csv")
                else pd.read_excel(uploaded)
            )
            st.dataframe(frame.head(50), hide_index=True, width="stretch")
            selected_columns: list[str] = []
            if bulk_mode == "Overwrite selected columns":
                if "NO MATRIK" not in frame.columns:
                    st.error("The uploaded file must contain NO MATRIK.")
                else:
                    technical = {"id", "created_at", "updated_at", "NO MATRIK"}
                    available_columns = [
                        column for column in frame.columns
                        if column not in technical
                        and (data.empty or column in data.columns)
                    ]
                    ignored = [
                        column for column in frame.columns
                        if column not in technical and column not in available_columns
                    ]
                    if ignored:
                        st.warning(
                            "Columns not found in the selected dataset will be ignored: "
                            + ", ".join(ignored)
                        )
                    selected_columns = st.multiselect(
                        "Columns to overwrite",
                        available_columns,
                        default=available_columns,
                    )
            confirmation = st.text_input("Type OVERWRITE to continue")
            disabled = (
                confirmation != "OVERWRITE"
                or (
                    bulk_mode == "Overwrite selected columns"
                    and ("NO MATRIK" not in frame.columns or not selected_columns)
                )
            )
            if st.button("Run bulk overwrite", disabled=disabled):
                try:
                    client = db()
                    if bulk_mode == "Replace entire dataset":
                        client.table(target).delete().neq(
                            "NO MATRIK", "__never__"
                        ).execute()
                        ok, message = upsert_rows(target, frame)
                        (st.success if ok else st.error)(message)
                    else:
                        existing_matrics = (
                            set(data["NO MATRIK"].dropna().astype(str))
                            if "NO MATRIK" in data.columns else set()
                        )
                        updated = 0
                        skipped = 0
                        for _, row in frame.iterrows():
                            matric = row.get("NO MATRIK")
                            if pd.isna(matric) or not str(matric).strip():
                                skipped += 1
                                continue
                            matric_value = str(matric).strip()
                            if existing_matrics and matric_value not in existing_matrics:
                                skipped += 1
                                continue
                            changes = {
                                column: (
                                    None
                                    if pd.isna(row[column])
                                    else (
                                        row[column].item()
                                        if hasattr(row[column], "item")
                                        else row[column]
                                    )
                                )
                                for column in selected_columns
                            }
                            result = client.table(target).update(changes).eq(
                                "NO MATRIK", matric_value
                            ).execute()
                            if result.data:
                                updated += 1
                            else:
                                skipped += 1
                        st.success(
                            f"{updated} records updated across "
                            f"{len(selected_columns)} selected columns."
                        )
                        if skipped:
                            st.warning(
                                f"{skipped} rows were skipped because NO MATRIK "
                                "was blank or not found in the selected dataset."
                            )
                except Exception as exc:
                    st.error(f"Overwrite stopped: {exc}")
    if not live:
        st.caption("The live dataset is unavailable or empty.")


def analytics_page() -> None:
    heading("Learning intelligence", "Performance analytics", "Turn class activity into clear, practical teaching decisions.")
    merged, _ = merged_students()
    numeric = [c for c in ["assignment", "pb", "task", "xp"] if c in merged]
    left, right = st.columns([1.4, 1])
    with left:
        if numeric:
            long = merged.melt(value_vars=numeric[:3], var_name="Assessment", value_name="Score")
            fig = px.box(long, x="Assessment", y="Score", color="Assessment", color_discrete_sequence=["#1F5DAA", "#17365D", "#8A96A3"])
            fig.update_layout(height=360, showlegend=False, margin=dict(l=5,r=5,t=20,b=5), plot_bgcolor="white", paper_bgcolor="white")
            polish_chart(fig)
            st.plotly_chart(fig, use_container_width=True)
    with right:
        st.subheader("Teaching signals")
        st.info("Published result performance is strongest in the leading assessment component.")
        st.warning("4 learners may need support in project-based assessment.")
        st.success("Task completion improved by 7% over the last four weeks.")
    if "xp" in merged:
        board = merged.sort_values("xp", ascending=False).reset_index(drop=True)
        board.insert(0, "rank", board.index + 1)
        st.subheader("Leaderboard")
        keep = [c for c in ["rank", "name", "student_name", "student_id", "xp", "badge"] if c in board]
        st.dataframe(board[keep], use_container_width=True, hide_index=True)


def admin_dashboard() -> None:
    heading("Administration", "XPLMS control centre", "Manage people, data quality and platform health with confidence.")
    merged, live = merged_students()
    a, b, c, d = st.columns(4)
    with a: metric("Student records", f"{len(merged):,}", "Active learners")
    with b: metric("Admin users", "16", "Across 6 programmes")
    with c: metric("Data health", "98.4%", "Last sync 12 min ago")
    with d: metric("Platform status", "Healthy", "All services normal")
    st.subheader("Quick actions")
    cols = st.columns(3)
    cols[0].button("＋ Add student", use_container_width=True)
    cols[1].button("⇧ Import dataset", use_container_width=True)
    cols[2].button("◉ Manage access", use_container_width=True)
    if not live: st.warning("Supabase tables are not currently available to the app. Apply the included database schema and verify your project URL/key.")


def user_access_page() -> None:
    heading(
        "Administration",
        "User access",
        "Create and maintain XPLMS admin and student accounts.",
    )
    client = db()
    if not client:
        st.error("The server connection is unavailable.")
        return
    users = pd.DataFrame(
        client.table("app_users")
        .select('id,username,email,full_name,role,"NO MATRIK",active,must_change_password,password_changed_at,last_login_at,created_at')
        .order("full_name")
        .execute()
        .data
    )
    create_tab, manage_tab, bulk_tab = st.tabs(
        ["Add user", "Manage users", "Bulk import"]
    )
    with create_tab:
        role = st.segmented_control("Access level", ["Admin", "Student"], default="Student")
        students, _ = fetch_table("stud_background", DEMO_STUDENTS)
        matric = None
        suggested_name = ""
        can_create = True
        if role == "Student":
            matric_col = "NO MATRIK" if "NO MATRIK" in students.columns else "student_id"
            name_col = "NAMA PELAJAR" if "NAMA PELAJAR" in students.columns else "name"
            assigned = set(users.get("NO MATRIK", pd.Series(dtype=str)).dropna().astype(str))
            available = students[~students[matric_col].astype(str).isin(assigned)].copy()
            options = available[matric_col].astype(str).tolist()
            if not options:
                st.info("Every student record is already linked to an account.")
                can_create = False
            else:
                name_map = dict(zip(available[matric_col].astype(str), available[name_col].astype(str)))
                matric = st.selectbox("Student record", options, format_func=lambda value: f"{name_map[value]} · {value}")
                suggested_name = name_map[matric]
        with st.form("create_app_user", clear_on_submit=True):
            full_name = st.text_input("Full name", value=suggested_name)
            email = st.text_input(
                "Admin email address" if role == "Admin" else "Student email address (optional)"
            )
            if role == "Student":
                st.info("The student's default username and password are both their NO MATRIK.")
                temp_password = matric or ""
            else:
                temp_password = st.text_input("Temporary password", type="password")
            if st.form_submit_button(
                "Create user", type="primary", disabled=not can_create
            ):
                if not full_name.strip() or (role == "Admin" and "@" not in email):
                    st.error("Enter a valid name and Admin email.")
                elif role == "Admin" and not valid_password(temp_password):
                    st.error("Use at least 10 characters with uppercase, lowercase and a number.")
                else:
                    try:
                        login_email = (
                            email.strip().lower()
                            if email.strip()
                            else f"{matric.lower()}@student.xplms.local"
                        )
                        client.table("app_users").insert({
                            "email": login_email,
                            "username": login_email if role == "Admin" else matric,
                            "full_name": full_name.strip(),
                            "role": role.lower(),
                            "NO MATRIK": matric if role == "Student" else None,
                            "password_hash": hash_password(temp_password),
                            "active": True,
                            "must_change_password": role == "Admin",
                        }).execute()
                        st.success(f"{role} account created. Share the temporary password securely.")
                    except Exception as exc:
                        st.error(f"User could not be created: {exc}")
    with manage_tab:
        if users.empty:
            st.info("No users found.")
            return
        visible_users = users.drop(columns=["id"], errors="ignore").copy()
        visible_users["password"] = visible_users.apply(
            lambda row: (
                str(row.get("NO MATRIK"))
                if (
                    str(row.get("role", "")).lower() == "student"
                    and pd.notna(row.get("NO MATRIK"))
                    and pd.isna(row.get("password_changed_at"))
                )
                else "PROTECTED (HASHED)"
            ),
            axis=1,
        )
        st.dataframe(
            visible_users,
            hide_index=True,
            width="stretch",
        )
        st.caption(
            "Only an unchanged default student password can be shown because it "
            "equals NO MATRIK. Changed passwords are securely hashed and cannot "
            "be recovered; use Reset password to replace one."
        )
        user_ids = users["id"].astype(str).tolist()
        label_map = {
            str(row["id"]): f"{row['full_name']} · {row['username']} · {row['role']}"
            for _, row in users.iterrows()
        }
        selected = st.selectbox("Select account", user_ids, format_func=lambda value: label_map[value])
        selected_row = users[users["id"].astype(str) == selected].iloc[0]
        c1, c2 = st.columns(2)
        active = c1.toggle("Account active", value=bool(selected_row["active"]))
        if c1.button("Save account status"):
            if selected == str(st.session_state.user["id"]) and not active:
                st.error("You cannot deactivate your own active session.")
            else:
                client.table("app_users").update({"active": active}).eq("id", selected).execute()
                st.success("Account status updated.")
        new_password = c2.text_input("New temporary password", type="password")
        if c2.button("Reset password"):
            selected_role = str(selected_row.get("role", "")).lower()
            password_is_valid = (
                len(new_password) >= 8
                if selected_role == "student"
                else valid_password(new_password)
            )
            if not password_is_valid:
                st.error(
                    "Use at least 8 characters for a student password. Admin "
                    "passwords require at least 10 characters with uppercase, "
                    "lowercase and a number."
                )
            else:
                client.table("app_users").update({
                    "password_hash": hash_password(new_password),
                    "must_change_password": True,
                    "password_changed_at": datetime.now().isoformat(),
                }).eq("id", selected).execute()
                st.success("Temporary password set. The user must change it at next login.")
        st.divider()
        remove_confirmed = st.checkbox(
            "Permanently remove this account",
            key="remove_user_confirmed",
        )
        if st.button("Remove user", disabled=not remove_confirmed):
            if selected == str(st.session_state.user["id"]):
                st.error("You cannot remove your own active account.")
            else:
                client.table("app_users").delete().eq("id", selected).execute()
                st.success("User account removed.")
                st.rerun()
    with bulk_tab:
        st.write(
            "Upload CSV or Excel with: full_name, role, NO MATRIK, email and "
            "temporary_password. Student email/password may be blank; NO MATRIK "
            "becomes both the default username and password."
        )
        template = pd.DataFrame([
            {
                "full_name": "Admin Name",
                "role": "admin",
                "NO MATRIK": "",
                "email": "admin@institution.edu",
                "temporary_password": "ChangeMe123",
            },
            {
                "full_name": "Student Name",
                "role": "student",
                "NO MATRIK": "MS00000000",
                "email": "",
                "temporary_password": "",
            },
        ])
        st.download_button(
            "Download CSV template",
            data=template.to_csv(index=False).encode("utf-8-sig"),
            file_name="xplms_user_access_template.csv",
            mime="text/csv",
            key="user_access_csv_template",
        )
        upload = st.file_uploader(
            "User access file",
            type=["csv", "xlsx", "xls"],
            key="user_access_bulk_file",
        )
        if upload:
            frame = (
                pd.read_csv(upload)
                if upload.name.lower().endswith(".csv")
                else pd.read_excel(upload)
            )
            st.dataframe(frame.head(50), hide_index=True, width="stretch")
            required = {"full_name", "role", "NO MATRIK"}
            missing = required - set(frame.columns)
            if missing:
                st.error(f"Missing columns: {', '.join(sorted(missing))}")
            elif st.button("Import user access", type="primary"):
                try:
                    records = []
                    for _, row in frame.iterrows():
                        role_value = str(row["role"]).strip().lower()
                        if role_value not in {"admin", "student"}:
                            raise ValueError(f"Invalid role: {row['role']}")
                        matric = row.get("NO MATRIK")
                        matric_value = str(matric).strip() if pd.notna(matric) else ""
                        raw_email = row.get("email", "")
                        email_value = (
                            ""
                            if pd.isna(raw_email)
                            or str(raw_email).strip().lower() in {"", "none", "nan", "null"}
                            else str(raw_email).strip().lower()
                        )
                        password = (
                            matric_value
                            if role_value == "student"
                            else str(row.get("temporary_password", ""))
                        )
                        if role_value == "student" and not matric_value:
                            raise ValueError("Every student account requires NO MATRIK.")
                        if role_value == "admin" and ("@" not in email_value or not valid_password(password)):
                            raise ValueError(
                                f"Admin {row['full_name']} requires a valid email and strong temporary password."
                            )
                        login_email = (
                            email_value
                            if email_value
                            else f"{matric_value.lower()}@student.xplms.local"
                        )
                        records.append({
                            "email": login_email,
                            "username": login_email if role_value == "admin" else matric_value,
                            "full_name": str(row["full_name"]).strip(),
                            "role": role_value,
                            "NO MATRIK": (
                                matric_value
                                if role_value == "student"
                                else None
                            ),
                            "password_hash": hash_password(password),
                            "active": True,
                            "must_change_password": role_value == "admin",
                        })
                    inserted = 0
                    updated = 0
                    for record in records:
                        if record["role"] == "student":
                            existing = (
                                client.table("app_users")
                                .select("id")
                                .eq("NO MATRIK", record["NO MATRIK"])
                                .limit(1)
                                .execute()
                                .data
                            )
                        else:
                            existing = (
                                client.table("app_users")
                                .select("id")
                                .ilike("email", record["email"])
                                .limit(1)
                                .execute()
                                .data
                            )
                        if existing:
                            client.table("app_users").update(record).eq(
                                "id", existing[0]["id"]
                            ).execute()
                            updated += 1
                        else:
                            client.table("app_users").insert(record).execute()
                            inserted += 1
                    st.success(
                        f"{len(records)} user accounts processed: "
                        f"{inserted} added, {updated} updated."
                    )
                except Exception as exc:
                    st.error(f"User import failed: {exc}")


def change_password_page() -> None:
    heading("", "Set a new password")
    st.write("Your administrator issued a temporary password. Replace it before continuing.")
    with st.form("change_password"):
        password = st.text_input("New password", type="password")
        confirm = st.text_input("Confirm new password", type="password")
        if st.form_submit_button("Update password", type="primary"):
            if password != confirm:
                st.error("The passwords do not match.")
            elif not valid_password(password):
                st.error("Use at least 10 characters with uppercase, lowercase and a number.")
            else:
                client = db()
                client.table("app_users").update({
                    "password_hash": hash_password(password),
                    "must_change_password": False,
                    "password_changed_at": datetime.now().isoformat(),
                }).eq("id", st.session_state.user["id"]).execute()
                st.session_state.user["must_change_password"] = False
                st.success("Password updated.")
                st.rerun()


def import_page() -> None:
    heading("Data operations", "Bulk import", "Review changes before replacing or updating a dataset.")
    target = st.selectbox("Target dataset", ["stud_background", "stud_progress"])
    uploaded = st.file_uploader("Upload CSV or Excel", type=["csv", "xlsx", "xls"])
    if uploaded:
        try:
            frame = pd.read_csv(uploaded) if uploaded.name.lower().endswith(".csv") else pd.read_excel(uploaded)
            st.success(f"{len(frame)} records ready · {len(frame.columns)} columns")
            st.dataframe(frame.head(50), use_container_width=True, hide_index=True)
            mode = st.radio("Import method", ["Update matching records", "Replace entire dataset"], horizontal=True)
            confirm = st.checkbox("I have reviewed this preview and want to continue")
            if st.button("Import records", type="primary", disabled=not confirm):
                client = db()
                if mode.startswith("Replace") and client:
                    try:
                        client.table(target).delete().neq("id", "__never__").execute()
                    except Exception as exc:
                        st.error(f"Replace was stopped safely: {exc}")
                        return
                ok, message = upsert_rows(target, frame)
                (st.success if ok else st.error)(message)
        except Exception as exc:
            st.error(f"Could not read this file: {exc}")
    st.download_button("Download student CSV template", DEMO_STUDENTS.head(0).to_csv(index=False), "student_import_template.csv", "text/csv")


def simple_page(title: str, copy: str) -> None:
    heading("XPLMS", title, copy)
    st.info("This workspace is ready for your institution-specific configuration.")


def main() -> None:
    demo_role = os.getenv("XPLMS_DEMO_ROLE", "").title()
    if "user" not in st.session_state and demo_role in {"Student", "Admin"}:
        st.session_state.user = {
            "name": f"Demo {demo_role}",
            "email": "demo@xplms.edu",
            "role": demo_role,
            "id": "demo",
        }
    if "user" not in st.session_state:
        login()
        return
    user = st.session_state.user
    if user.get("must_change_password"):
        st.session_state.header_identity = (
            f"{user.get('name', 'USER')} | {user.get('role', 'USER')}"
        ).upper()
        change_password_page()
        return
    actual_role = user.get("role", "Student")
    role = (
        st.session_state.get("admin_preview_role", "Admin")
        if actual_role == "Admin" and not is_demo()
        else actual_role
    )
    if role not in {"Student", "Admin"}:
        role = "Student"
    if role == "Student":
        display_name, student_class = student_identity(user)
        st.session_state.header_identity = (
            f"{display_name} | {student_class}"
        ).upper()
    else:
        display_name = user["name"]
        st.session_state.header_identity = f"{display_name} | ADMIN".upper()
    page = sidebar(role, display_name)
    page = mobile_navigation(role, page)
    if role == "Student":
        {
            "Results": progress_page,
            "Materials": materials_page,
            "Quiz": quiz_page,
            "XP Journey": my_xp_page,
            "SOP": sop_page,
            "Leaderboard": student_leaderboard_page,
            "Profile": profile_page,
            "Request XP": request_xp_page,
        }[page]()
    else:
        {
            "Student record": admin_student_records_page,
            "Analysis background": analysis_background_page,
            "Analysis results": analysis_progress_page,
            "Analysis XP": analysis_xp_page,
            "CRUD": admin_crud_page,
            "Award XP": award_xp_page,
            "SOP": sop_page,
            "User access": user_access_page,
            "Material": lambda: materials_page(True),
            "Quiz": admin_quiz_page,
        }[page]()


if __name__ == "__main__":
    main()
